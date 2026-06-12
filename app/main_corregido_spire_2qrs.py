from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Cm
from io import BytesIO
import unicodedata
import re
import os
import json
from datetime import datetime
from dotenv import load_dotenv
from openai import OpenAI
from copy import deepcopy
from typing import List, Tuple
import qrcode
import threading
import tempfile

# -------------------------------------------------
# MERGE PROFESIONAL OPCIONAL PARA UNIVERSIDAD_2QRS
# -------------------------------------------------
# python-pptx no fusiona de forma segura presentaciones con masters/layouts distintos.
# Para UNIVERSIDAD_2QRS, donde sí necesitas mezclar DIPLOMADO, CURSO, etc.
# dentro de un solo PPTX, usamos Spire.Presentation solo para el merge final.
try:
    from spire.presentation import Presentation as SpirePresentation
    from spire.presentation.common import FileFormat
except Exception:
    SpirePresentation = None
    FileFormat = None

# -------------------------------------------------
# CONFIGURACIÓN GENERAL
# -------------------------------------------------

load_dotenv()
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
client = OpenAI(api_key=OPENAI_API_KEY)

app = FastAPI(title="Certificados API", version="2.3.0")

# Modelos disponibles (carpetas dentro de app/templates)
# IMPORTANTE: Deben coincidir con los valores enviados desde el frontend:
# - INSTITUTO
# - UNIVERSIDAD_2QRS
# - UNIVERSIDAD_AZUL
# - COLEGIO_ABOGADOS_CALLAO
MODELO_FOLDER_MAP = {
    "INSTITUTO": "instituto",
    "UNIVERSIDAD_2QRS": "universidad_2qrs",
    "UNIVERSIDAD_AZUL": "universidad_azul",
    "COLEGIO_ABOGADOS_CALLAO": "colegio_de_abogados_del_callao", 
    "COLEGIO_DE_PROFESORES_DEL_PERU": "colegio_de_profesores_del_peru",
}

# Modelos que usan formato de fecha largo (dd de Mes del yyyy)
MODELOS_FECHA_LARGA = {
    "INSTITUTO",
    "UNIVERSIDAD_AZUL",
    "COLEGIO_ABOGADOS_CALLAO", 
    "COLEGIO_DE_PROFESORES_DEL_PERU",
}


# Nombre de archivo de plantilla por tipo (dentro de cada carpeta de modelo)
TEMPLATE_FILENAME_MAP = {
    "DIPLOMADO": "plantilla_diplomado.pptx",
    "PROGRAMA DE ESPECIALIZACIÓN": "plantilla_programa_de_especializacion.pptx",
    "CURSO": "plantilla_curso.pptx",
    "CURSO_DE_CAPACITACION": "plantilla_curso_de_capacitacion.pptx",
    "CURSO_DE_ACTUALIZACION": "plantilla_curso_de_actualizacion.pptx",
}

# Nº de módulos por tipo
MODULOS_COUNT = {
    "DIPLOMADO": 8,
    "PROGRAMA DE ESPECIALIZACIÓN": 8,
    "CURSO": 5,
    "CURSO_DE_CAPACITACION": 5,
    "CURSO_DE_ACTUALIZACION": 5,
}

# -------------------------------------------------
# JSON PERSISTENTE DE MÓDULOS
# -------------------------------------------------
# JSON base incluido en el repositorio GitHub.
# Este archivo sirve como base fija después de que descargues el JSON del volumen
# y lo subas manualmente a tu repositorio.
MODULOS_BASE_JSON_PATH = os.getenv(
    "MODULOS_BASE_JSON_PATH",
    os.path.join("app", "modulos_base.json")
)

# JSON persistente de Railway Volume.
# En Railway, RAILWAY_VOLUME_MOUNT_PATH lo crea automáticamente el volumen.
# En local, usará app/data/modulos_cache.json.
MODULOS_VOLUME_DIR = os.getenv(
    "RAILWAY_VOLUME_MOUNT_PATH",
    os.path.join("app", "data")
)
MODULOS_VOLUME_JSON_PATH = os.getenv(
    "MODULOS_VOLUME_JSON_PATH",
    os.path.join(MODULOS_VOLUME_DIR, "modulos_cache.json")
)

MODULOS_JSON_LOCK = threading.Lock()


# -------------------------------------------------
# UTILIDADES
# -------------------------------------------------

def safe_filename(text: str) -> str:
    text = unicodedata.normalize("NFKD", text).encode("ascii", "ignore").decode("ascii")
    text = re.sub(r"[^a-zA-Z0-9_-]", "_", text)
    return text


def format_date_ddmmyyyy(date_str: str) -> str:
    try:
        return datetime.strptime(date_str, "%Y-%m-%d").strftime("%d/%m/%Y")
    except ValueError:
        return date_str


def format_date_long_es(date_str: str) -> str:
    try:
        date_obj = datetime.strptime(date_str, "%Y-%m-%d")
        meses = [
            "Enero", "Febrero", "Marzo", "Abril", "Mayo", "Junio",
            "Julio", "Agosto", "Septiembre", "Octubre", "Noviembre", "Diciembre"
        ]

        # ✅ día con cero a la izquierda
        dia = f"{date_obj.day:02d}"

        return f"{dia} de {meses[date_obj.month - 1]} del {date_obj.year}"
    except ValueError:
        return date_str

def format_date_range_long_es(fecha_inicio: str, fecha_fin: str) -> tuple[str, str]:
    """
    Devuelve (fecha_inicio_formateada, fecha_fin_formateada)
    aplicando la regla:
    - Si ambos años son iguales → el año solo se muestra en la fecha final
    - Si son distintos → cada fecha muestra su año
    """
    try:
        inicio = datetime.strptime(fecha_inicio, "%Y-%m-%d")
        fin = datetime.strptime(fecha_fin, "%Y-%m-%d")

        meses = [
            "Enero", "Febrero", "Marzo", "Abril", "Mayo", "Junio",
            "Julio", "Agosto", "Septiembre", "Octubre", "Noviembre", "Diciembre"
        ]

        dia_inicio = f"{inicio.day:02d}"
        dia_fin = f"{fin.day:02d}"

        if inicio.year == fin.year:
            fecha_inicio_str = f"{dia_inicio} de {meses[inicio.month - 1]}"
            fecha_fin_str = f"{dia_fin} de {meses[fin.month - 1]} del {fin.year}"
        else:
            fecha_inicio_str = f"{dia_inicio} de {meses[inicio.month - 1]} del {inicio.year}"
            fecha_fin_str = f"{dia_fin} de {meses[fin.month - 1]} del {fin.year}"

        return fecha_inicio_str, fecha_fin_str

    except ValueError:
        # fallback seguro
        return format_date_long_es(fecha_inicio), format_date_long_es(fecha_fin)


def format_two_digits_number(value: int) -> str:
    """
    Formatea números enteros a dos dígitos.
    Ej: 3 -> 03, 12 -> 12
    """
    try:
        return f"{int(value):02d}"
    except (ValueError, TypeError):
        return str(value)


def format_two_digits_float(value: float) -> str:
    """
    Formatea números decimales manteniendo decimales,
    pero con parte entera a dos dígitos.
    Ej: 3 -> 03
        3.5 -> 03.5
        12 -> 12
        12.25 -> 12.25
    """
    try:
        entero = int(value)
        decimal = value - entero

        if decimal == 0:
            return f"{entero:02d}"

        # Eliminar ceros innecesarios en decimales
        decimal_str = str(round(decimal, 2)).lstrip("0")
        return f"{entero:02d}{decimal_str}"
    except (ValueError, TypeError):
        return str(value)


def calcular_horas_por_modulo(total_horas: int, cantidad_modulos: int) -> str:
    """
    Calcula horas por módulo:
    total_horas / cantidad_modulos
    Devuelve string sin decimales si es entero, o con 2 decimales si no.
    """
    if cantidad_modulos <= 0:
        return "0"

    valor = total_horas / cantidad_modulos

    # Si es entero, no mostrar decimales
    if valor.is_integer():
        return str(int(valor))

    # Si no, mostrar hasta 2 decimales
    return f"{valor:.2f}"


def nombre_completo_capitalizado(nombres: str, apellidos: str) -> str:
    texto = f"{nombres} {apellidos}".strip().lower()
    return " ".join(p.capitalize() for p in texto.split())

# cambio 1
def obtener_iniciales_apellidos(apellidos: str) -> str:
    """
    Obtiene las iniciales de los dos primeros apellidos.
    Ej:
    'Chávez Guzmán' -> 'CG'
    'Carlos de la Cruz' -> 'CD'
    """
    partes = (apellidos or "").strip().split()

    if not partes:
        return ""

    iniciales = "".join(parte[0] for parte in partes[:2] if parte)

    return iniciales.upper()


def modelo_con_mayuscula_inicial(tipo_modelo: str) -> str:
    texto = tipo_modelo.strip().lower().split()
    if not texto:
        return ""
    texto[0] = texto[0].capitalize()
    return " ".join(texto)


def resolve_template_path(modelo_certificado: str, tipo_modelo: str) -> str:
    modelo_key = (modelo_certificado or "").upper().strip()
    tipo_key = (tipo_modelo or "").upper().strip()

    if modelo_key not in MODELO_FOLDER_MAP:
        raise HTTPException(
            status_code=400,
            detail=f"Modelo de certificado no soportado: {modelo_key}"
        )

    if tipo_key not in TEMPLATE_FILENAME_MAP:
        raise HTTPException(
            status_code=400,
            detail=f"Tipo de certificado no soportado: {tipo_key}"
        )

    folder = MODELO_FOLDER_MAP[modelo_key]
    filename = TEMPLATE_FILENAME_MAP[tipo_key]

    return os.path.join("app", "templates", folder, filename)

##agregué
def normalize_text_for_url(text: str) -> str:
    text = unicodedata.normalize("NFKD", text)
    text = text.encode("ascii", "ignore").decode("ascii")
    text = re.sub(r"[^a-zA-Z0-9]", "", text)
    return text.lower()


def build_qr_url(nombres: str, apellidos: str, tema: str) -> str:
    primer_nombre = nombres.strip().split()[0]
    primer_apellido = apellidos.strip().split()[0]

    base = (
        normalize_text_for_url(primer_nombre)
        + normalize_text_for_url(primer_apellido)
        + normalize_text_for_url(tema)
    )

    return f"https://especializacionvirtual.com/certificados/{base}.pdf"


def generate_qr_image(url: str) -> BytesIO:
    qr = qrcode.QRCode(
        version=4,
        error_correction=qrcode.constants.ERROR_CORRECT_Q,
        box_size=10,
        border=1,
    )
    qr.add_data(url)
    qr.make(fit=True)

    img = qr.make_image(fill_color="black", back_color="white")

    buffer = BytesIO()
    img.save(buffer, format="PNG")
    buffer.seek(0)
    return buffer
##fin agregué

def insert_qr_at_placeholder_in_slide(slide, qr_stream: BytesIO, qr_size_cm: float = 2.78) -> int:
    """
    Inserta el QR únicamente dentro del slide recibido.

    Antes se recorría toda la presentación desde el inicio cada vez que se quería
    insertar un QR. Eso podía asociar el QR de un participante con una diapositiva
    incorrecta cuando el PPTX final tenía varios certificados combinados.
    """
    QR_SIZE = Cm(qr_size_cm)
    inserted_count = 0

    # Guardamos los bytes para poder insertar el mismo QR en más de un placeholder
    # sin depender de la posición actual del stream.
    qr_bytes = qr_stream.getvalue()

    for shape in list(slide.shapes):
        if not shape.has_text_frame:
            continue

        if "{{QR_CODE}}" not in shape.text_frame.text:
            continue

        left = shape.left
        top = shape.top

        # Eliminamos el placeholder para evitar cajas vacías encima del QR.
        element = shape._element
        element.getparent().remove(element)

        slide.shapes.add_picture(
            BytesIO(qr_bytes),
            left=left,
            top=top,
            width=QR_SIZE,
            height=QR_SIZE,
        )

        inserted_count += 1

    return inserted_count


# Compatibilidad: se mantiene la función anterior por si alguna parte externa
# todavía la invoca. Internamente el endpoint usa la versión slide por slide.
def insert_qr_at_placeholder(prs: Presentation, qr_stream: BytesIO, qr_size_cm: float = 2.78):
    for slide in prs.slides:
        inserted = insert_qr_at_placeholder_in_slide(slide, qr_stream, qr_size_cm)
        if inserted > 0:
            return


##DISTRIBUIR HORAS POR MÓDULO
def distribuir_horas_por_modulo(total_horas: int, cantidad_modulos: int) -> List[int]:
    """
    Distribuye las horas de forma progresiva y variada,
    asegurando que la suma final sea EXACTA.
    """
    if cantidad_modulos <= 0:
        return []

    pesos = []
    incremento = 0.1
    for i in range(cantidad_modulos):
        peso = 1 + (i * incremento)
        pesos.append(peso)
    suma_pesos = sum(pesos)
    
    horas = [
        int((peso / suma_pesos) * total_horas)
        for peso in pesos
    ]

    # Ajuste para asegurar suma exacta
    diferencia = total_horas - sum(horas)

    # Repartir diferencia empezando por el último módulo
    i = cantidad_modulos - 1
    while diferencia > 0:
        horas[i] += 1
        diferencia -= 1
        i -= 1
        if i < 0:
            i = cantidad_modulos - 1

    return horas


# -------------------------------------------------
# AJUSTE DE TABLA SEGÚN LONGITUD DEL TEMA
# -------------------------------------------------

EMU_PER_PT = 12700

LABEL_TIPO_MAP = {
    "DIPLOMADO": "DIPLOMADO",
    "PROGRAMA DE ESPECIALIZACIÓN": "PROGRAMA DE ESPECIALIZACIÓN",
    "CURSO": "CURSO",
    "CURSO_DE_CAPACITACION": "CURSO DE CAPACITACIÓN",
    "CURSO_DE_ACTUALIZACION": "CURSO DE ACTUALIZACIÓN",
}


def find_shape_by_name(slide, name: str):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def estimate_chars_per_line(shape_width_emu: int, font_size_pt: float) -> int:
    """
    Estima cuántos caracteres entran por línea en un shape.
    No mide el render exacto de PowerPoint; usa una aproximación segura.
    """
    if not font_size_pt or font_size_pt <= 0:
        font_size_pt = 12.0

    width_pt = shape_width_emu / EMU_PER_PT
    avg_char_pt = font_size_pt * 0.52
    chars_per_line = int(width_pt / avg_char_pt)

    safety_factor = 0.85
    chars_per_line = int(chars_per_line * safety_factor)

    return max(10, chars_per_line)


def wrap_by_words(text: str, max_chars: int) -> list[str]:
    """
    Simula un ajuste de línea por palabras.
    Si una palabra supera el máximo, la divide para evitar desbordes.
    """
    words = text.split(" ")
    lines = []
    current = ""

    for word in words:
        candidate = word if current == "" else current + " " + word

        if len(candidate) <= max_chars:
            current = candidate
            continue

        if current:
            lines.append(current)
            current = ""

        while len(word) > max_chars:
            lines.append(word[:max_chars])
            word = word[max_chars:]

        current = word

    if current:
        lines.append(current)

    return lines


def ajustar_tabla_certificado_estudios_generico(
    prs: Presentation,
    label: str,
    tema: str,
    *,
    shape_tema_name: str = "PH_TEMA",
    shape_tabla_name: str = "PH_TABLA",
    gap_min=Cm(0.25),
    max_extra_lines: int = 3,
):
    """
    Mueve hacia abajo PH_TABLA cuando el texto de PH_TEMA ocupa más de una línea.
    Requiere que los shapes estén nombrados como PH_TEMA y PH_TABLA en la plantilla.
    """
    label = (label or "").strip().upper()
    tema = (tema or "").strip()
    texto_linea = f"{label}: {tema}".strip()

    for slide in prs.slides:
        ph_tema = find_shape_by_name(slide, shape_tema_name)
        ph_tabla = find_shape_by_name(slide, shape_tabla_name)

        if not ph_tema or not ph_tabla:
            continue

        font_pt = 12.0

        try:
            text_frame = ph_tema.text_frame
            for paragraph in text_frame.paragraphs:
                paragraph_text = "".join(run.text for run in paragraph.runs).upper()

                if label and label in paragraph_text:
                    for run in paragraph.runs:
                        if run.font and run.font.size:
                            font_pt = run.font.size.pt
                            raise StopIteration
        except StopIteration:
            pass
        except Exception:
            pass

        max_chars = estimate_chars_per_line(ph_tema.width, font_pt)
        lines = wrap_by_words(texto_linea, max_chars)
        line_count = len(lines)

        if line_count <= 1:
            return

        if line_count > max_extra_lines:
            line_count = max_extra_lines

        line_height_emu = int((font_pt * 0.65) * EMU_PER_PT)
        extra_lines = line_count - 1
        delta = extra_lines * line_height_emu

        push_extra = Cm(0.00)
        if line_count >= 3:
            push_extra = Cm(0.25)

        ph_tabla.top = ph_tabla.top + delta + int(gap_min) + int(push_extra)
        return



# -------------------------------------------------
# CACHE PERSISTENTE DE MÓDULOS EN JSON
# -------------------------------------------------

def normalizar_clave_tema(texto: str) -> str:
    """
    Normaliza el tema para usarlo como clave estable del JSON.
    Ej:
    'Didáctica de la Comunicación en Educación Secundaria'
    -> 'didactica-de-la-comunicacion-en-educacion-secundaria'
    """
    texto = unicodedata.normalize("NFKD", texto or "")
    texto = texto.encode("ascii", "ignore").decode("ascii")
    texto = texto.lower().strip()
    texto = re.sub(r"[^a-z0-9]+", "-", texto)
    texto = re.sub(r"-+", "-", texto).strip("-")
    return texto or "sin-tema"


def grupo_modulos_por_tipo(tipo: str) -> str:
    """
    Agrupa los tipos por cantidad de módulos.
    DIPLOMADO y PROGRAMA DE ESPECIALIZACIÓN comparten los mismos 8 módulos.
    CURSO, CURSO DE CAPACITACIÓN y CURSO DE ACTUALIZACIÓN comparten 5 módulos.
    """
    tipo_key = (tipo or "").upper().strip()

    if tipo_key not in MODULOS_COUNT:
        raise ValueError("Tipo no soportado para módulos")

    cantidad = MODULOS_COUNT[tipo_key]
    return f"{cantidad}_MODULOS"


def construir_clave_modulos(tipo: str, tema: str) -> str:
    grupo = grupo_modulos_por_tipo(tipo)
    tema_key = normalizar_clave_tema(tema)
    return f"{grupo}::{tema_key}"


def empty_modulos_json() -> dict:
    return {
        "version": 1,
        "items": {}
    }


def leer_json_modulos(path: str) -> dict:
    """
    Lee un JSON de módulos.
    Si no existe o está dañado, devuelve una estructura vacía.
    """
    try:
        if not path or not os.path.exists(path):
            return empty_modulos_json()

        with open(path, "r", encoding="utf-8") as f:
            data = json.load(f)

        if not isinstance(data, dict):
            return empty_modulos_json()

        if "items" not in data or not isinstance(data.get("items"), dict):
            # Soporte por si alguna vez el archivo fue guardado como dict directo.
            return {
                "version": 1,
                "items": data
            }

        return data

    except Exception:
        return empty_modulos_json()


def extraer_modulos_de_registro(registro) -> list[str] | None:
    """
    Soporta dos formatos:
    1) {"modulos": [...]}
    2) [...] directamente
    """
    if isinstance(registro, list):
        modulos = registro
    elif isinstance(registro, dict):
        modulos = registro.get("modulos")
    else:
        return None

    if not isinstance(modulos, list):
        return None

    return [str(m).upper().strip() for m in modulos if str(m).strip()]


def buscar_modulos_en_json(path: str, cache_key: str, cantidad_esperada: int) -> list[str] | None:
    data = leer_json_modulos(path)
    items = data.get("items", {})

    registro = items.get(cache_key)
    modulos = extraer_modulos_de_registro(registro)

    if not modulos:
        return None

    if len(modulos) != cantidad_esperada:
        return None

    return modulos


def guardar_modulos_en_volume_json(
    cache_key: str,
    tipo: str,
    tema: str,
    modulos: list[str],
) -> None:
    """
    Guarda módulos nuevos en el JSON persistente del Railway Volume.
    No guarda en RAM y no modifica el JSON base del repositorio.
    """
    with MODULOS_JSON_LOCK:
        os.makedirs(os.path.dirname(MODULOS_VOLUME_JSON_PATH), exist_ok=True)

        data = leer_json_modulos(MODULOS_VOLUME_JSON_PATH)
        items = data.setdefault("items", {})

        now = datetime.now().isoformat(timespec="seconds")

        registro_anterior = items.get(cache_key)
        created_at = now
        if isinstance(registro_anterior, dict) and registro_anterior.get("created_at"):
            created_at = registro_anterior["created_at"]

        items[cache_key] = {
            "grupo": grupo_modulos_por_tipo(tipo),
            "tipo_referencia": (tipo or "").upper().strip(),
            "tema_original": (tema or "").strip(),
            "modulos": [str(m).upper().strip() for m in modulos],
            "created_at": created_at,
            "updated_at": now,
        }

        tmp_path = f"{MODULOS_VOLUME_JSON_PATH}.tmp"

        with open(tmp_path, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)

        os.replace(tmp_path, MODULOS_VOLUME_JSON_PATH)


def obtener_modulos_desde_json(tipo: str, tema: str) -> list[str] | None:
    tipo_key = (tipo or "").upper().strip()
    if tipo_key not in MODULOS_COUNT:
        raise ValueError("Tipo no soportado para módulos")

    cantidad = MODULOS_COUNT[tipo_key]
    cache_key = construir_clave_modulos(tipo_key, tema)

    # 1. Primero Railway Volume: es el más reciente.
    modulos_volume = buscar_modulos_en_json(
        MODULOS_VOLUME_JSON_PATH,
        cache_key,
        cantidad
    )
    if modulos_volume:
        return modulos_volume

    # 2. Luego JSON base del repositorio GitHub.
    modulos_base = buscar_modulos_en_json(
        MODULOS_BASE_JSON_PATH,
        cache_key,
        cantidad
    )
    if modulos_base:
        return modulos_base

    return None


# -------------------------------------------------
# OPENAI – GENERACIÓN DE MÓDULOS DINÁMICA
# -------------------------------------------------

def build_prompt(tema: str, count: int) -> str:
    return f"""
Devuelve EXCLUSIVAMENTE un JSON válido con este formato exacto:

{{
  "modulos": [
    {', '.join(['"string"' for _ in range(count)])}
  ]
}}

REGLAS OBLIGATORIAS:
- Exactamente {count} módulos
- SOLO títulos (NO descripciones)
- Máximo 12 palabras por módulo
- En español
- NO numeración
- NO texto fuera del JSON
- NO explicaciones

Certificado: {tema}
""".strip()


def obtener_modulos_por_tema(tipo: str, tema: str) -> list[str]:
    tipo_key = (tipo or "").upper().strip()

    if tipo_key not in MODULOS_COUNT:
        raise ValueError("Tipo no soportado para módulos")

    count = MODULOS_COUNT[tipo_key]
    cache_key = construir_clave_modulos(tipo_key, tema)

    # 1. Buscar primero en JSON persistente:
    #    - Railway Volume
    #    - JSON base del repositorio GitHub
    modulos_json = obtener_modulos_desde_json(tipo_key, tema)
    if modulos_json:
        return modulos_json

    # 2. Si no existe en ningún JSON, recién usar OpenAI.
    try:
        response = client.responses.create(
            model="gpt-5-mini",
            input=[
                {"role": "system", "content": "Eres un asistente académico extremadamente estricto."},
                {"role": "user", "content": build_prompt(tema, count)}
            ]
        )

        data = json.loads(response.output_text)
        modulos = data.get("modulos", [])

        if not isinstance(modulos, list) or len(modulos) != count:
            raise ValueError("Cantidad de módulos inválida")

        modulos = [str(m).upper().strip() for m in modulos]

        # 3. Guardar solo resultados válidos de OpenAI en Railway Volume.
        guardar_modulos_en_volume_json(
            cache_key=cache_key,
            tipo=tipo_key,
            tema=tema,
            modulos=modulos,
        )

        return modulos

    except Exception:
        # No guardar fallback en JSON para no contaminar el cache persistente.
        return [f"MÓDULO {i+1}" for i in range(count)]

# -------------------------------------------------
# CORS
# -------------------------------------------------

app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "http://localhost:4200",
        "http://127.0.0.1:4200",
        "https://ipdefrontendcertificados.vercel.app", 
    ],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


# -------------------------------------------------
# DTO
# -------------------------------------------------

class DiplomaRequest(BaseModel):
    modeloCertificado: str
    tipoModelo: str
    nombres: str
    apellidos: str
    temaDiplomado: str
    fechaInicio: str
    fechaFin: str
    horasAcademicas: int
    creditosAcademicos: int
    folioNumero: str
    fechaEmision: str
    codigoEstudiante: str = ""
    ciudad: str = ""


class BatchRequest(BaseModel):
    items: List[DiplomaRequest]


# -------------------------------------------------
# REEMPLAZO DE TEXTO (PRESERVA FUENTES)
# -------------------------------------------------

def replace_in_text_frame_preserve_font(text_frame, mapping):
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if not run.text:
                continue
            for k, v in mapping.items():
                if k in run.text:
                    run.text = run.text.replace(k, v)


def replace_in_shape(shape, mapping):
    if shape.has_text_frame:
        replace_in_text_frame_preserve_font(shape.text_frame, mapping)

    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                replace_in_text_frame_preserve_font(cell.text_frame, mapping)

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shape.shapes:
            replace_in_shape(s, mapping)


def replace_placeholders(prs, mapping):
    for slide in prs.slides:
        for shape in slide.shapes:
            replace_in_shape(shape, mapping)

        if slide.has_notes_slide:
            for shape in slide.notes_slide.shapes:
                replace_in_shape(shape, mapping)

    for layout in prs.slide_layouts:
        for shape in layout.shapes:
            replace_in_shape(shape, mapping)

    for master in prs.slide_masters:
        for shape in master.shapes:
            replace_in_shape(shape, mapping)


# -------------------------------------------------
# MERGE PPTX
# -------------------------------------------------

UNIVERSIDAD_2QRS_MODEL_KEY = "UNIVERSIDAD_2QRS"


def _replace_rids_in_element(el, rid_map: dict[str, str]):
    """
    Reemplaza ids de relación dentro de un fragmento XML copiado.
    Esto es necesario para que imágenes y recursos relacionados apunten al slide destino.
    """
    for e in el.iter():
        if not hasattr(e, "attrib"):
            continue
        for attr_key, attr_val in list(e.attrib.items()):
            if attr_val in rid_map:
                e.attrib[attr_key] = rid_map[attr_val]


def _element_has_placeholder(el) -> bool:
    """
    Detecta si un elemento XML es un placeholder.
    Los placeholders heredados de layout/master no deben materializarse como shapes reales,
    porque pueden duplicar cajas de texto o alterar el diseño final.
    """
    for e in el.iter():
        if str(e.tag).lower().endswith("}ph"):
            return True
    return False


def _iter_copyable_sp_tree_children(sp_tree, *, skip_placeholders: bool):
    """
    Itera los hijos realmente copiables de un spTree.
    No copia propiedades internas del grupo raíz del slide/layout/master.
    """
    for child in list(sp_tree):
        tag = child.tag.lower()

        # Estos nodos son propiedades internas del grupo raíz.
        # Copiarlos genera XML inválido o comportamientos visuales extraños.
        if tag.endswith("nvgrpsppr") or tag.endswith("grpsppr"):
            continue

        if skip_placeholders and _element_has_placeholder(child):
            continue

        yield child


def _clone_child_with_relationships(dest_slide, src_part, child):
    """
    Clona un shape XML y recrea sus relaciones en el slide destino.
    Aplica especialmente a imágenes, fondos, logotipos y otros recursos embebidos.
    """
    cloned_child = deepcopy(child)

    rid_map = {}
    for rId, rel in src_part.rels.items():
        try:
            new_rId = dest_slide.part.relate_to(
                rel._target,
                rel.reltype,
                is_external=rel.is_external,
            )
            rid_map[rId] = new_rId
        except Exception:
            # Si una relación no aplica al shape clonado, se omite sin romper el PPTX.
            continue

    _replace_rids_in_element(cloned_child, rid_map)
    return cloned_child


def _copy_sp_tree_children_to_slide(dest_slide, src_part, src_sp_tree, *, skip_placeholders: bool):
    """
    Copia shapes desde un slide, layout o master hacia un slide destino.
    Se usa en UNIVERSIDAD_2QRS para materializar el diseño visual heredado
    antes de ocultar el master/layout del destino.
    """
    dest_sp_tree = dest_slide.shapes._spTree

    for child in _iter_copyable_sp_tree_children(
        src_sp_tree,
        skip_placeholders=skip_placeholders,
    ):
        cloned_child = _clone_child_with_relationships(dest_slide, src_part, child)
        dest_sp_tree.insert_element_before(cloned_child, "p:extLst")


def _remove_slide_level_shapes(slide):
    """
    Elimina shapes creados automáticamente al agregar una diapositiva nueva.
    No elimina elementos heredados del master, por eso en 2QRS también se usa
    showMasterSp=0 y una materialización controlada del diseño fuente.
    """
    for shape in list(slide.shapes):
        element = shape._element
        element.getparent().remove(element)


def _hide_inherited_master_shapes(slide):
    """
    Evita que la diapositiva clonada herede shapes del master del PPTX destino.

    Esto es clave para UNIVERSIDAD_2QRS: cuando se mezclan DIPLOMADO, CURSO,
    CERTIFICADO DE ESTUDIOS, etc. dentro del mismo modelo, usar el master/layout
    del primer template puede superponer o agrandar el fondo del jaguar/escudo.
    """
    slide._element.set("showMasterSp", "0")


def _count_non_placeholder_shapes_in_layout(layout) -> int:
    """
    Cuenta shapes visuales no-placeholder en un layout.
    Sirve para elegir el layout más limpio del destino en UNIVERSIDAD_2QRS.
    """
    count = 0
    for _ in _iter_copyable_sp_tree_children(
        layout.shapes._spTree,
        skip_placeholders=True,
    ):
        count += 1
    return count


def _get_cleanest_layout(prs: Presentation):
    """
    Devuelve el layout con menos shapes visuales propios.
    Es más seguro que asumir que slide_layouts[6] siempre está realmente vacío.
    """
    layouts = list(prs.slide_layouts)
    if not layouts:
        raise ValueError("La presentación no tiene layouts disponibles")

    return min(layouts, key=_count_non_placeholder_shapes_in_layout)


def clone_slide_into(dest_prs: Presentation, src_slide):
    """
    Clonación original para los modelos que ya funcionan correctamente.
    No se cambia su comportamiento para evitar afectar INSTITUTO, UNIVERSIDAD_AZUL,
    COLEGIO_ABOGADOS_CALLAO o COLEGIO_DE_PROFESORES_DEL_PERU.
    """
    blank_layout = dest_prs.slide_layouts[6]
    new_slide = dest_prs.slides.add_slide(blank_layout)

    spTree = new_slide.shapes._spTree
    src_spTree = src_slide.shapes._spTree

    for child in list(src_spTree):
        tag = child.tag.lower()

        if tag.endswith("nvgrpsppr") or tag.endswith("grpsppr"):
            continue

        spTree.insert_element_before(deepcopy(child), "p:extLst")

    rid_map = {}
    for rId, rel in src_slide.part.rels.items():
        try:
            new_rId = new_slide.part.relate_to(
                rel._target,
                rel.reltype,
                is_external=rel.is_external,
            )
            rid_map[rId] = new_rId
        except Exception:
            continue

    _replace_rids_in_element(new_slide._element, rid_map)


def clone_slide_into_universidad_2qrs(dest_prs: Presentation, src_slide):
    """
    Clonación especial SOLO para UNIVERSIDAD_2QRS.

    Motivo:
    - En 2QRS, algunas piezas visuales del diseño pueden vivir en master/layout.
    - Al mezclar distintos tipoModelo en el mismo modeloCertificado, el slide clonado
      puede heredar el master/layout del primer PPTX del lote.
    - Eso produce fondos gigantes, marcas de agua duplicadas o elementos superpuestos.

    Solución:
    1. Crear un slide con el layout más limpio del destino.
    2. Ocultar shapes heredados del master del destino.
    3. Materializar primero el master y layout del slide fuente.
    4. Copiar encima los shapes reales del slide fuente.
    """
    clean_layout = _get_cleanest_layout(dest_prs)
    new_slide = dest_prs.slides.add_slide(clean_layout)

    _hide_inherited_master_shapes(new_slide)
    _remove_slide_level_shapes(new_slide)

    # 1) Materializar shapes visuales del master fuente, sin placeholders.
    try:
        src_master = src_slide.slide_layout.slide_master
        _copy_sp_tree_children_to_slide(
            new_slide,
            src_master.part,
            src_master.shapes._spTree,
            skip_placeholders=True,
        )
    except Exception:
        pass

    # 2) Materializar shapes visuales del layout fuente, sin placeholders.
    try:
        src_layout = src_slide.slide_layout
        _copy_sp_tree_children_to_slide(
            new_slide,
            src_layout.part,
            src_layout.shapes._spTree,
            skip_placeholders=True,
        )
    except Exception:
        pass

    # 3) Copiar shapes propios del slide fuente, incluyendo placeholders ya rellenados.
    _copy_sp_tree_children_to_slide(
        new_slide,
        src_slide.part,
        src_slide.shapes._spTree,
        skip_placeholders=False,
    )

    return new_slide


def merge_presentations(
    presentations: List[Presentation],
    modelo_certificado: str = "",
) -> Presentation:
    if not presentations:
        raise ValueError("No hay presentaciones para unir")

    modelo_key = (modelo_certificado or "").upper().strip()

    # Usar la primera presentación como base preserva su tamaño de página,
    # tema y configuración general del archivo.
    dest = presentations[0]

    for prs in presentations[1:]:
        for slide in prs.slides:
            if modelo_key == UNIVERSIDAD_2QRS_MODEL_KEY:
                clone_slide_into_universidad_2qrs(dest, slide)
            else:
                clone_slide_into(dest, slide)

    return dest


# -------------------------------------------------
# GENERACIÓN DE PPTX POR ITEM  ✅ RESTAURADA
# -------------------------------------------------

def generar_presentacion_por_item(item: DiplomaRequest) -> Presentation:
    tipo = item.tipoModelo.upper().strip()
    modelo_cert = item.modeloCertificado.upper().strip()

    template_path = resolve_template_path(modelo_cert, tipo)

    if not os.path.exists(template_path):
        raise HTTPException(
            status_code=500,
            detail=f"No existe la plantilla: {template_path}"
        )

    prs = Presentation(template_path)
    
    # Determinar formato de fecha según modelo de certificado
    usar_fecha_larga = modelo_cert in MODELOS_FECHA_LARGA

    if usar_fecha_larga:
        fecha_inicio, fecha_fin = format_date_range_long_es(
            item.fechaInicio,
            item.fechaFin
        )
        fecha_emision_larga = format_date_long_es(item.fechaEmision)
        fecha_emision_corta = format_date_ddmmyyyy(item.fechaEmision)
        # fecha_emision = format_date_long_es(item.fechaEmision)
    else:
        fecha_inicio = format_date_ddmmyyyy(item.fechaInicio)
        fecha_fin = format_date_ddmmyyyy(item.fechaFin)
        fecha_emision_larga = format_date_long_es(item.fechaEmision)
        fecha_emision_corta = format_date_ddmmyyyy(item.fechaEmision)
        # fecha_emision = format_date_ddmmyyyy(item.fechaEmision)
    
    nombre_posterior = nombre_completo_capitalizado(item.nombres, item.apellidos)
    modulos = obtener_modulos_por_tema(tipo, item.temaDiplomado)
    
    ##horas de módulos correctamente distribuidas
    cantidad_modulos = MODULOS_COUNT[tipo]

    horas_por_modulo_global = calcular_horas_por_modulo(
        item.horasAcademicas,
        cantidad_modulos
    )
    
    horas_por_modulo = distribuir_horas_por_modulo(
        item.horasAcademicas,
        cantidad_modulos
    )

    mapping = {
        "{{MODELO}}": tipo,
        "{{MODELO_MINUSCULA}}": modelo_con_mayuscula_inicial(tipo),

        "{{NOMBRES}}": item.nombres.upper(),
        "{{APELLIDOS}}": item.apellidos.upper(),
        "{{TEMA_DIPLOMADO}}": item.temaDiplomado.upper(),

        "{{NOMBRE_COMPLETO_MINUSCULA}}": nombre_posterior,

        "{{FECHA_INICIO}}": fecha_inicio,
        "{{FECHA_FIN}}": fecha_fin,

        "{{HORAS_ACADEMICAS}}": format_two_digits_number(item.horasAcademicas),
        "{{CREDITOS_ACADEMICOS}}": format_two_digits_float(item.creditosAcademicos),
        
        "{{HORAS_MODULO}}": horas_por_modulo_global,

        "{{FOLIO_NUMERO}}": item.folioNumero,

        "{{FECHA_EMISION_LARGA}}": fecha_emision_larga,
        "{{FECHA_EMISION_CORTA}}": fecha_emision_corta,

        "{{CODE_STUDENT}}": (item.codigoEstudiante or "").strip().lower(),
        "{{INICIALES}}": obtener_iniciales_apellidos(item.apellidos),
        "{{CIUDAD}}": (item.ciudad or "").strip(),
    }

    for i, m in enumerate(modulos, start=1):
        mapping[f"{{{{MODULO_{i}}}}}"] = m
    
    ##Agregar horas por módulo al mapping       
    for i, h in enumerate(horas_por_modulo, start=1):
        mapping[f"{{{{HORAS_MODULO_{i}}}}}"] = str(h)


    replace_placeholders(prs, mapping)

    label_tabla = LABEL_TIPO_MAP.get(tipo)
    if label_tabla:
        ajustar_tabla_certificado_estudios_generico(
            prs,
            label=label_tabla,
            tema=item.temaDiplomado.upper(),
        )

    return prs



# -------------------------------------------------
# MERGE SEGURO PARA UNIVERSIDAD_2QRS USANDO SPIRE
# -------------------------------------------------

def presentation_to_bytes(prs: Presentation) -> bytes:
    """
    Convierte una Presentation de python-pptx a bytes PPTX.
    """
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


def merge_pptx_bytes_with_spire(pptx_files: list[bytes]) -> bytes:
    """
    Fusiona PPTX completos preservando diseños, masters, layouts e imágenes.

    Esta función se usa SOLO para UNIVERSIDAD_2QRS, porque ese modelo necesita
    mezclar distintos tipoModelo en el mismo archivo final y python-pptx no hace
    ese merge de forma segura cuando hay masters/layouts diferentes.
    """
    if not pptx_files:
        raise ValueError("No hay PPTX para fusionar")

    if SpirePresentation is None or FileFormat is None:
        raise HTTPException(
            status_code=500,
            detail=(
                "Falta instalar Spire.Presentation. Agrega 'Spire.Presentation' "
                "a requirements.txt y redeploya en Railway."
            ),
        )

    with tempfile.TemporaryDirectory() as tmpdir:
        input_paths: list[str] = []

        for index, pptx_bytes in enumerate(pptx_files, start=1):
            path = os.path.join(tmpdir, f"input_{index}.pptx")
            with open(path, "wb") as f:
                f.write(pptx_bytes)
            input_paths.append(path)

        output_path = os.path.join(tmpdir, "merged_2qrs.pptx")

        base = SpirePresentation()
        base.LoadFromFile(input_paths[0])

        try:
            for path in input_paths[1:]:
                src = SpirePresentation()
                src.LoadFromFile(path)

                try:
                    # AppendBySlide preserva el diseño original de la diapositiva fuente.
                    for slide in src.Slides:
                        base.Slides.AppendBySlide(slide)
                finally:
                    src.Dispose()

            base.SaveToFile(output_path, FileFormat.Pptx2016)
        finally:
            base.Dispose()

        with open(output_path, "rb") as f:
            return f.read()


def generar_pptx_2qrs_con_qr_como_bytes(item: DiplomaRequest) -> bytes:
    """
    Genera el PPTX individual de un item UNIVERSIDAD_2QRS, inserta su QR
    dentro de su propio archivo y devuelve bytes listos para fusionar.
    """
    prs = generar_presentacion_por_item(item)

    qr_url = build_qr_url(
        item.nombres,
        item.apellidos,
        item.temaDiplomado,
    )
    qr_image = generate_qr_image(qr_url)

    for slide in prs.slides:
        insert_qr_at_placeholder_in_slide(
            slide,
            qr_image,
            qr_size_cm=2.78,
        )

    return presentation_to_bytes(prs)

# -------------------------------------------------
# ENDPOINTS
# -------------------------------------------------

def obtener_modelo_unico_del_lote(items: List[DiplomaRequest]) -> str:
    """
    Valida que el lote pertenezca a un solo modeloCertificado.

    Sí se permite mezclar distintos tipoModelo, por ejemplo:
    - UNIVERSIDAD_2QRS + DIPLOMADO
    - UNIVERSIDAD_2QRS + CURSO
    - UNIVERSIDAD_2QRS + CURSO_DE_CAPACITACION

    Lo que no debe ocurrir es mezclar UNIVERSIDAD_2QRS con UNIVERSIDAD_AZUL,
    INSTITUTO u otros modelos dentro del mismo PPTX.
    """
    modelo_base = items[0].modeloCertificado.upper().strip()

    for index, item in enumerate(items, start=1):
        modelo_actual = item.modeloCertificado.upper().strip()
        if modelo_actual != modelo_base:
            raise HTTPException(
                status_code=400,
                detail=(
                    "No mezcles distintos modeloCertificado en un mismo lote. "
                    f"El primer item usa {modelo_base}, pero el item {index} usa {modelo_actual}. "
                    "Puedes mezclar tipoModelo, pero no modeloCertificado."
                ),
            )

    return modelo_base


@app.get("/health")
def health():
    return {"status": "ok"}


@app.post("/api/diplomas")
def generate_pptx_batch(
    payload: BatchRequest,
):
    if not payload.items:
        raise HTTPException(status_code=400, detail="items no puede estar vacío")

    modelo_lote = obtener_modelo_unico_del_lote(payload.items)

    # -------------------------------------------------
    # CASO ESPECIAL: UNIVERSIDAD_2QRS
    # -------------------------------------------------
    # Este modelo es el único que se deformaba al mezclar distintos tipoModelo.
    # La causa era el merge manual con python-pptx: copiaba shapes, pero no
    # preservaba de forma segura masters/layouts/relaciones internas.
    #
    # Solución: cada certificado se genera completo por separado, se inserta su QR
    # en su propio PPTX y al final se fusionan los PPTX completos con Spire.
    if modelo_lote == UNIVERSIDAD_2QRS_MODEL_KEY:
        pptx_files: list[bytes] = []

        for item in payload.items:
            pptx_files.append(generar_pptx_2qrs_con_qr_como_bytes(item))

        merged_bytes = merge_pptx_bytes_with_spire(pptx_files)
        output = BytesIO(merged_bytes)
        output.seek(0)

        filename = f"CERTIFICADOS_{int(datetime.now().timestamp())}.pptx"

        return StreamingResponse(
            output,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f'attachment; filename="{filename}"'}
        )

    # -------------------------------------------------
    # FLUJO ORIGINAL PARA LOS DEMÁS MODELOS
    # -------------------------------------------------
    # No se toca porque indicaste que los otros modelos funcionan correctamente.
    presentations: List[Presentation] = []
    slide_counts: List[int] = []

    for item in payload.items:
        prs = generar_presentacion_por_item(item)
        presentations.append(prs)
        slide_counts.append(len(prs.slides))

    merged = merge_presentations(
        presentations,
        modelo_certificado=modelo_lote,
    )

    merged_slides = list(merged.slides)
    cursor = 0

    for item, slide_count in zip(payload.items, slide_counts):
        qr_url = build_qr_url(
            item.nombres,
            item.apellidos,
            item.temaDiplomado,
        )
        qr_image = generate_qr_image(qr_url)

        qr_size_cm = 2.0 if item.modeloCertificado.upper().strip() == "UNIVERSIDAD_AZUL" else 2.78

        slides_del_item = merged_slides[cursor: cursor + slide_count]
        for slide in slides_del_item:
            insert_qr_at_placeholder_in_slide(
                slide,
                qr_image,
                qr_size_cm=qr_size_cm,
            )

        cursor += slide_count

    output = BytesIO()
    merged.save(output)
    output.seek(0)

    filename = f"CERTIFICADOS_{int(datetime.now().timestamp())}.pptx"

    return StreamingResponse(
        output,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'}
    )
