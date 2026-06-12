from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Cm
from io import BytesIO
import unicodedata
import re
import os
import json
from datetime import datetime
from dotenv import load_dotenv
from openai import OpenAI
from copy import deepcopy
from typing import List, Tuple
import qrcode
import threading
import zipfile
import posixpath
from xml.etree import ElementTree as ET

# -------------------------------------------------
# CONFIGURACIÓN GENERAL
# -------------------------------------------------

load_dotenv()
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
client = OpenAI(api_key=OPENAI_API_KEY)

app = FastAPI(title="Certificados API", version="2.3.0")

# Modelos disponibles (carpetas dentro de app/templates)
# IMPORTANTE: Deben coincidir con los valores enviados desde el frontend:
# - INSTITUTO
# - UNIVERSIDAD_2QRS
# - UNIVERSIDAD_AZUL
# - COLEGIO_ABOGADOS_CALLAO
MODELO_FOLDER_MAP = {
    "INSTITUTO": "instituto",
    "UNIVERSIDAD_2QRS": "universidad_2qrs",
    "UNIVERSIDAD_AZUL": "universidad_azul",
    "COLEGIO_ABOGADOS_CALLAO": "colegio_de_abogados_del_callao", 
    "COLEGIO_DE_PROFESORES_DEL_PERU": "colegio_de_profesores_del_peru",
}

# Modelos que usan formato de fecha largo (dd de Mes del yyyy)
MODELOS_FECHA_LARGA = {
    "INSTITUTO",
    "UNIVERSIDAD_AZUL",
    "COLEGIO_ABOGADOS_CALLAO", 
    "COLEGIO_DE_PROFESORES_DEL_PERU",
}


# Nombre de archivo de plantilla por tipo (dentro de cada carpeta de modelo)
TEMPLATE_FILENAME_MAP = {
    "DIPLOMADO": "plantilla_diplomado.pptx",
    "PROGRAMA DE ESPECIALIZACIÓN": "plantilla_programa_de_especializacion.pptx",
    "CURSO": "plantilla_curso.pptx",
    "CURSO_DE_CAPACITACION": "plantilla_curso_de_capacitacion.pptx",
    "CURSO_DE_ACTUALIZACION": "plantilla_curso_de_actualizacion.pptx",
}

# Nº de módulos por tipo
MODULOS_COUNT = {
    "DIPLOMADO": 8,
    "PROGRAMA DE ESPECIALIZACIÓN": 8,
    "CURSO": 5,
    "CURSO_DE_CAPACITACION": 5,
    "CURSO_DE_ACTUALIZACION": 5,
}

# -------------------------------------------------
# JSON PERSISTENTE DE MÓDULOS
# -------------------------------------------------
# JSON base incluido en el repositorio GitHub.
# Este archivo sirve como base fija después de que descargues el JSON del volumen
# y lo subas manualmente a tu repositorio.
MODULOS_BASE_JSON_PATH = os.getenv(
    "MODULOS_BASE_JSON_PATH",
    os.path.join("app", "modulos_base.json")
)

# JSON persistente de Railway Volume.
# En Railway, RAILWAY_VOLUME_MOUNT_PATH lo crea automáticamente el volumen.
# En local, usará app/data/modulos_cache.json.
MODULOS_VOLUME_DIR = os.getenv(
    "RAILWAY_VOLUME_MOUNT_PATH",
    os.path.join("app", "data")
)
MODULOS_VOLUME_JSON_PATH = os.getenv(
    "MODULOS_VOLUME_JSON_PATH",
    os.path.join(MODULOS_VOLUME_DIR, "modulos_cache.json")
)

MODULOS_JSON_LOCK = threading.Lock()


# -------------------------------------------------
# UTILIDADES
# -------------------------------------------------

def safe_filename(text: str) -> str:
    text = unicodedata.normalize("NFKD", text).encode("ascii", "ignore").decode("ascii")
    text = re.sub(r"[^a-zA-Z0-9_-]", "_", text)
    return text


def format_date_ddmmyyyy(date_str: str) -> str:
    try:
        return datetime.strptime(date_str, "%Y-%m-%d").strftime("%d/%m/%Y")
    except ValueError:
        return date_str


def format_date_long_es(date_str: str) -> str:
    try:
        date_obj = datetime.strptime(date_str, "%Y-%m-%d")
        meses = [
            "Enero", "Febrero", "Marzo", "Abril", "Mayo", "Junio",
            "Julio", "Agosto", "Septiembre", "Octubre", "Noviembre", "Diciembre"
        ]

        # ✅ día con cero a la izquierda
        dia = f"{date_obj.day:02d}"

        return f"{dia} de {meses[date_obj.month - 1]} del {date_obj.year}"
    except ValueError:
        return date_str

def format_date_range_long_es(fecha_inicio: str, fecha_fin: str) -> tuple[str, str]:
    """
    Devuelve (fecha_inicio_formateada, fecha_fin_formateada)
    aplicando la regla:
    - Si ambos años son iguales → el año solo se muestra en la fecha final
    - Si son distintos → cada fecha muestra su año
    """
    try:
        inicio = datetime.strptime(fecha_inicio, "%Y-%m-%d")
        fin = datetime.strptime(fecha_fin, "%Y-%m-%d")

        meses = [
            "Enero", "Febrero", "Marzo", "Abril", "Mayo", "Junio",
            "Julio", "Agosto", "Septiembre", "Octubre", "Noviembre", "Diciembre"
        ]

        dia_inicio = f"{inicio.day:02d}"
        dia_fin = f"{fin.day:02d}"

        if inicio.year == fin.year:
            fecha_inicio_str = f"{dia_inicio} de {meses[inicio.month - 1]}"
            fecha_fin_str = f"{dia_fin} de {meses[fin.month - 1]} del {fin.year}"
        else:
            fecha_inicio_str = f"{dia_inicio} de {meses[inicio.month - 1]} del {inicio.year}"
            fecha_fin_str = f"{dia_fin} de {meses[fin.month - 1]} del {fin.year}"

        return fecha_inicio_str, fecha_fin_str

    except ValueError:
        # fallback seguro
        return format_date_long_es(fecha_inicio), format_date_long_es(fecha_fin)


def format_two_digits_number(value: int) -> str:
    """
    Formatea números enteros a dos dígitos.
    Ej: 3 -> 03, 12 -> 12
    """
    try:
        return f"{int(value):02d}"
    except (ValueError, TypeError):
        return str(value)


def format_two_digits_float(value: float) -> str:
    """
    Formatea números decimales manteniendo decimales,
    pero con parte entera a dos dígitos.
    Ej: 3 -> 03
        3.5 -> 03.5
        12 -> 12
        12.25 -> 12.25
    """
    try:
        entero = int(value)
        decimal = value - entero

        if decimal == 0:
            return f"{entero:02d}"

        # Eliminar ceros innecesarios en decimales
        decimal_str = str(round(decimal, 2)).lstrip("0")
        return f"{entero:02d}{decimal_str}"
    except (ValueError, TypeError):
        return str(value)


def calcular_horas_por_modulo(total_horas: int, cantidad_modulos: int) -> str:
    """
    Calcula horas por módulo:
    total_horas / cantidad_modulos
    Devuelve string sin decimales si es entero, o con 2 decimales si no.
    """
    if cantidad_modulos <= 0:
        return "0"

    valor = total_horas / cantidad_modulos

    # Si es entero, no mostrar decimales
    if valor.is_integer():
        return str(int(valor))

    # Si no, mostrar hasta 2 decimales
    return f"{valor:.2f}"


def nombre_completo_capitalizado(nombres: str, apellidos: str) -> str:
    texto = f"{nombres} {apellidos}".strip().lower()
    return " ".join(p.capitalize() for p in texto.split())

# cambio 1
def obtener_iniciales_apellidos(apellidos: str) -> str:
    """
    Obtiene las iniciales de los dos primeros apellidos.
    Ej:
    'Chávez Guzmán' -> 'CG'
    'Carlos de la Cruz' -> 'CD'
    """
    partes = (apellidos or "").strip().split()

    if not partes:
        return ""

    iniciales = "".join(parte[0] for parte in partes[:2] if parte)

    return iniciales.upper()


def modelo_con_mayuscula_inicial(tipo_modelo: str) -> str:
    texto = tipo_modelo.strip().lower().split()
    if not texto:
        return ""
    texto[0] = texto[0].capitalize()
    return " ".join(texto)


def resolve_template_path(modelo_certificado: str, tipo_modelo: str) -> str:
    modelo_key = (modelo_certificado or "").upper().strip()
    tipo_key = (tipo_modelo or "").upper().strip()

    if modelo_key not in MODELO_FOLDER_MAP:
        raise HTTPException(
            status_code=400,
            detail=f"Modelo de certificado no soportado: {modelo_key}"
        )

    if tipo_key not in TEMPLATE_FILENAME_MAP:
        raise HTTPException(
            status_code=400,
            detail=f"Tipo de certificado no soportado: {tipo_key}"
        )

    folder = MODELO_FOLDER_MAP[modelo_key]
    filename = TEMPLATE_FILENAME_MAP[tipo_key]

    return os.path.join("app", "templates", folder, filename)

##agregué
def normalize_text_for_url(text: str) -> str:
    text = unicodedata.normalize("NFKD", text)
    text = text.encode("ascii", "ignore").decode("ascii")
    text = re.sub(r"[^a-zA-Z0-9]", "", text)
    return text.lower()


def build_qr_url(nombres: str, apellidos: str, tema: str) -> str:
    primer_nombre = nombres.strip().split()[0]
    primer_apellido = apellidos.strip().split()[0]

    base = (
        normalize_text_for_url(primer_nombre)
        + normalize_text_for_url(primer_apellido)
        + normalize_text_for_url(tema)
    )

    return f"https://especializacionvirtual.com/certificados/{base}.pdf"


def generate_qr_image(url: str) -> BytesIO:
    qr = qrcode.QRCode(
        version=4,
        error_correction=qrcode.constants.ERROR_CORRECT_Q,
        box_size=10,
        border=1,
    )
    qr.add_data(url)
    qr.make(fit=True)

    img = qr.make_image(fill_color="black", back_color="white")

    buffer = BytesIO()
    img.save(buffer, format="PNG")
    buffer.seek(0)
    return buffer
##fin agregué

def insert_qr_at_placeholder_in_slide(slide, qr_stream: BytesIO, qr_size_cm: float = 2.78) -> int:
    """
    Inserta el QR únicamente dentro del slide recibido.

    Antes se recorría toda la presentación desde el inicio cada vez que se quería
    insertar un QR. Eso podía asociar el QR de un participante con una diapositiva
    incorrecta cuando el PPTX final tenía varios certificados combinados.
    """
    QR_SIZE = Cm(qr_size_cm)
    inserted_count = 0

    # Guardamos los bytes para poder insertar el mismo QR en más de un placeholder
    # sin depender de la posición actual del stream.
    qr_bytes = qr_stream.getvalue()

    for shape in list(slide.shapes):
        if not shape.has_text_frame:
            continue

        if "{{QR_CODE}}" not in shape.text_frame.text:
            continue

        left = shape.left
        top = shape.top

        # Eliminamos el placeholder para evitar cajas vacías encima del QR.
        element = shape._element
        element.getparent().remove(element)

        slide.shapes.add_picture(
            BytesIO(qr_bytes),
            left=left,
            top=top,
            width=QR_SIZE,
            height=QR_SIZE,
        )

        inserted_count += 1

    return inserted_count


# Compatibilidad: se mantiene la función anterior por si alguna parte externa
# todavía la invoca. Internamente el endpoint usa la versión slide por slide.
def insert_qr_at_placeholder(prs: Presentation, qr_stream: BytesIO, qr_size_cm: float = 2.78):
    for slide in prs.slides:
        inserted = insert_qr_at_placeholder_in_slide(slide, qr_stream, qr_size_cm)
        if inserted > 0:
            return


##DISTRIBUIR HORAS POR MÓDULO
def distribuir_horas_por_modulo(total_horas: int, cantidad_modulos: int) -> List[int]:
    """
    Distribuye las horas de forma progresiva y variada,
    asegurando que la suma final sea EXACTA.
    """
    if cantidad_modulos <= 0:
        return []

    pesos = []
    incremento = 0.1
    for i in range(cantidad_modulos):
        peso = 1 + (i * incremento)
        pesos.append(peso)
    suma_pesos = sum(pesos)
    
    horas = [
        int((peso / suma_pesos) * total_horas)
        for peso in pesos
    ]

    # Ajuste para asegurar suma exacta
    diferencia = total_horas - sum(horas)

    # Repartir diferencia empezando por el último módulo
    i = cantidad_modulos - 1
    while diferencia > 0:
        horas[i] += 1
        diferencia -= 1
        i -= 1
        if i < 0:
            i = cantidad_modulos - 1

    return horas


# -------------------------------------------------
# AJUSTE DE TABLA SEGÚN LONGITUD DEL TEMA
# -------------------------------------------------

EMU_PER_PT = 12700

LABEL_TIPO_MAP = {
    "DIPLOMADO": "DIPLOMADO",
    "PROGRAMA DE ESPECIALIZACIÓN": "PROGRAMA DE ESPECIALIZACIÓN",
    "CURSO": "CURSO",
    "CURSO_DE_CAPACITACION": "CURSO DE CAPACITACIÓN",
    "CURSO_DE_ACTUALIZACION": "CURSO DE ACTUALIZACIÓN",
}


def find_shape_by_name(slide, name: str):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def estimate_chars_per_line(shape_width_emu: int, font_size_pt: float) -> int:
    """
    Estima cuántos caracteres entran por línea en un shape.
    No mide el render exacto de PowerPoint; usa una aproximación segura.
    """
    if not font_size_pt or font_size_pt <= 0:
        font_size_pt = 12.0

    width_pt = shape_width_emu / EMU_PER_PT
    avg_char_pt = font_size_pt * 0.52
    chars_per_line = int(width_pt / avg_char_pt)

    safety_factor = 0.85
    chars_per_line = int(chars_per_line * safety_factor)

    return max(10, chars_per_line)


def wrap_by_words(text: str, max_chars: int) -> list[str]:
    """
    Simula un ajuste de línea por palabras.
    Si una palabra supera el máximo, la divide para evitar desbordes.
    """
    words = text.split(" ")
    lines = []
    current = ""

    for word in words:
        candidate = word if current == "" else current + " " + word

        if len(candidate) <= max_chars:
            current = candidate
            continue

        if current:
            lines.append(current)
            current = ""

        while len(word) > max_chars:
            lines.append(word[:max_chars])
            word = word[max_chars:]

        current = word

    if current:
        lines.append(current)

    return lines


def ajustar_tabla_certificado_estudios_generico(
    prs: Presentation,
    label: str,
    tema: str,
    *,
    shape_tema_name: str = "PH_TEMA",
    shape_tabla_name: str = "PH_TABLA",
    gap_min=Cm(0.25),
    max_extra_lines: int = 3,
):
    """
    Mueve hacia abajo PH_TABLA cuando el texto de PH_TEMA ocupa más de una línea.
    Requiere que los shapes estén nombrados como PH_TEMA y PH_TABLA en la plantilla.
    """
    label = (label or "").strip().upper()
    tema = (tema or "").strip()
    texto_linea = f"{label}: {tema}".strip()

    for slide in prs.slides:
        ph_tema = find_shape_by_name(slide, shape_tema_name)
        ph_tabla = find_shape_by_name(slide, shape_tabla_name)

        if not ph_tema or not ph_tabla:
            continue

        font_pt = 12.0

        try:
            text_frame = ph_tema.text_frame
            for paragraph in text_frame.paragraphs:
                paragraph_text = "".join(run.text for run in paragraph.runs).upper()

                if label and label in paragraph_text:
                    for run in paragraph.runs:
                        if run.font and run.font.size:
                            font_pt = run.font.size.pt
                            raise StopIteration
        except StopIteration:
            pass
        except Exception:
            pass

        max_chars = estimate_chars_per_line(ph_tema.width, font_pt)
        lines = wrap_by_words(texto_linea, max_chars)
        line_count = len(lines)

        if line_count <= 1:
            return

        if line_count > max_extra_lines:
            line_count = max_extra_lines

        line_height_emu = int((font_pt * 0.65) * EMU_PER_PT)
        extra_lines = line_count - 1
        delta = extra_lines * line_height_emu

        push_extra = Cm(0.00)
        if line_count >= 3:
            push_extra = Cm(0.25)

        ph_tabla.top = ph_tabla.top + delta + int(gap_min) + int(push_extra)
        return



# -------------------------------------------------
# CACHE PERSISTENTE DE MÓDULOS EN JSON
# -------------------------------------------------

def normalizar_clave_tema(texto: str) -> str:
    """
    Normaliza el tema para usarlo como clave estable del JSON.
    Ej:
    'Didáctica de la Comunicación en Educación Secundaria'
    -> 'didactica-de-la-comunicacion-en-educacion-secundaria'
    """
    texto = unicodedata.normalize("NFKD", texto or "")
    texto = texto.encode("ascii", "ignore").decode("ascii")
    texto = texto.lower().strip()
    texto = re.sub(r"[^a-z0-9]+", "-", texto)
    texto = re.sub(r"-+", "-", texto).strip("-")
    return texto or "sin-tema"


def grupo_modulos_por_tipo(tipo: str) -> str:
    """
    Agrupa los tipos por cantidad de módulos.
    DIPLOMADO y PROGRAMA DE ESPECIALIZACIÓN comparten los mismos 8 módulos.
    CURSO, CURSO DE CAPACITACIÓN y CURSO DE ACTUALIZACIÓN comparten 5 módulos.
    """
    tipo_key = (tipo or "").upper().strip()

    if tipo_key not in MODULOS_COUNT:
        raise ValueError("Tipo no soportado para módulos")

    cantidad = MODULOS_COUNT[tipo_key]
    return f"{cantidad}_MODULOS"


def construir_clave_modulos(tipo: str, tema: str) -> str:
    grupo = grupo_modulos_por_tipo(tipo)
    tema_key = normalizar_clave_tema(tema)
    return f"{grupo}::{tema_key}"


def empty_modulos_json() -> dict:
    return {
        "version": 1,
        "items": {}
    }


def leer_json_modulos(path: str) -> dict:
    """
    Lee un JSON de módulos.
    Si no existe o está dañado, devuelve una estructura vacía.
    """
    try:
        if not path or not os.path.exists(path):
            return empty_modulos_json()

        with open(path, "r", encoding="utf-8") as f:
            data = json.load(f)

        if not isinstance(data, dict):
            return empty_modulos_json()

        if "items" not in data or not isinstance(data.get("items"), dict):
            # Soporte por si alguna vez el archivo fue guardado como dict directo.
            return {
                "version": 1,
                "items": data
            }

        return data

    except Exception:
        return empty_modulos_json()


def extraer_modulos_de_registro(registro) -> list[str] | None:
    """
    Soporta dos formatos:
    1) {"modulos": [...]}
    2) [...] directamente
    """
    if isinstance(registro, list):
        modulos = registro
    elif isinstance(registro, dict):
        modulos = registro.get("modulos")
    else:
        return None

    if not isinstance(modulos, list):
        return None

    return [str(m).upper().strip() for m in modulos if str(m).strip()]


def buscar_modulos_en_json(path: str, cache_key: str, cantidad_esperada: int) -> list[str] | None:
    data = leer_json_modulos(path)
    items = data.get("items", {})

    registro = items.get(cache_key)
    modulos = extraer_modulos_de_registro(registro)

    if not modulos:
        return None

    if len(modulos) != cantidad_esperada:
        return None

    return modulos


def guardar_modulos_en_volume_json(
    cache_key: str,
    tipo: str,
    tema: str,
    modulos: list[str],
) -> None:
    """
    Guarda módulos nuevos en el JSON persistente del Railway Volume.
    No guarda en RAM y no modifica el JSON base del repositorio.
    """
    with MODULOS_JSON_LOCK:
        os.makedirs(os.path.dirname(MODULOS_VOLUME_JSON_PATH), exist_ok=True)

        data = leer_json_modulos(MODULOS_VOLUME_JSON_PATH)
        items = data.setdefault("items", {})

        now = datetime.now().isoformat(timespec="seconds")

        registro_anterior = items.get(cache_key)
        created_at = now
        if isinstance(registro_anterior, dict) and registro_anterior.get("created_at"):
            created_at = registro_anterior["created_at"]

        items[cache_key] = {
            "grupo": grupo_modulos_por_tipo(tipo),
            "tipo_referencia": (tipo or "").upper().strip(),
            "tema_original": (tema or "").strip(),
            "modulos": [str(m).upper().strip() for m in modulos],
            "created_at": created_at,
            "updated_at": now,
        }

        tmp_path = f"{MODULOS_VOLUME_JSON_PATH}.tmp"

        with open(tmp_path, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)

        os.replace(tmp_path, MODULOS_VOLUME_JSON_PATH)


def obtener_modulos_desde_json(tipo: str, tema: str) -> list[str] | None:
    tipo_key = (tipo or "").upper().strip()
    if tipo_key not in MODULOS_COUNT:
        raise ValueError("Tipo no soportado para módulos")

    cantidad = MODULOS_COUNT[tipo_key]
    cache_key = construir_clave_modulos(tipo_key, tema)

    # 1. Primero Railway Volume: es el más reciente.
    modulos_volume = buscar_modulos_en_json(
        MODULOS_VOLUME_JSON_PATH,
        cache_key,
        cantidad
    )
    if modulos_volume:
        return modulos_volume

    # 2. Luego JSON base del repositorio GitHub.
    modulos_base = buscar_modulos_en_json(
        MODULOS_BASE_JSON_PATH,
        cache_key,
        cantidad
    )
    if modulos_base:
        return modulos_base

    return None


# -------------------------------------------------
# OPENAI – GENERACIÓN DE MÓDULOS DINÁMICA
# -------------------------------------------------

def build_prompt(tema: str, count: int) -> str:
    return f"""
Devuelve EXCLUSIVAMENTE un JSON válido con este formato exacto:

{{
  "modulos": [
    {', '.join(['"string"' for _ in range(count)])}
  ]
}}

REGLAS OBLIGATORIAS:
- Exactamente {count} módulos
- SOLO títulos (NO descripciones)
- Máximo 12 palabras por módulo
- En español
- NO numeración
- NO texto fuera del JSON
- NO explicaciones

Certificado: {tema}
""".strip()


def obtener_modulos_por_tema(tipo: str, tema: str) -> list[str]:
    tipo_key = (tipo or "").upper().strip()

    if tipo_key not in MODULOS_COUNT:
        raise ValueError("Tipo no soportado para módulos")

    count = MODULOS_COUNT[tipo_key]
    cache_key = construir_clave_modulos(tipo_key, tema)

    # 1. Buscar primero en JSON persistente:
    #    - Railway Volume
    #    - JSON base del repositorio GitHub
    modulos_json = obtener_modulos_desde_json(tipo_key, tema)
    if modulos_json:
        return modulos_json

    # 2. Si no existe en ningún JSON, recién usar OpenAI.
    try:
        response = client.responses.create(
            model="gpt-5-mini",
            input=[
                {"role": "system", "content": "Eres un asistente académico extremadamente estricto."},
                {"role": "user", "content": build_prompt(tema, count)}
            ]
        )

        data = json.loads(response.output_text)
        modulos = data.get("modulos", [])

        if not isinstance(modulos, list) or len(modulos) != count:
            raise ValueError("Cantidad de módulos inválida")

        modulos = [str(m).upper().strip() for m in modulos]

        # 3. Guardar solo resultados válidos de OpenAI en Railway Volume.
        guardar_modulos_en_volume_json(
            cache_key=cache_key,
            tipo=tipo_key,
            tema=tema,
            modulos=modulos,
        )

        return modulos

    except Exception:
        # No guardar fallback en JSON para no contaminar el cache persistente.
        return [f"MÓDULO {i+1}" for i in range(count)]

# -------------------------------------------------
# CORS
# -------------------------------------------------

app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "http://localhost:4200",
        "http://127.0.0.1:4200",
        "https://ipdefrontendcertificados.vercel.app", 
    ],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


# -------------------------------------------------
# DTO
# -------------------------------------------------

class DiplomaRequest(BaseModel):
    modeloCertificado: str
    tipoModelo: str
    nombres: str
    apellidos: str
    temaDiplomado: str
    fechaInicio: str
    fechaFin: str
    horasAcademicas: int
    creditosAcademicos: int
    folioNumero: str
    fechaEmision: str
    codigoEstudiante: str = ""
    ciudad: str = ""


class BatchRequest(BaseModel):
    items: List[DiplomaRequest]


# -------------------------------------------------
# REEMPLAZO DE TEXTO (PRESERVA FUENTES)
# -------------------------------------------------

def replace_in_text_frame_preserve_font(text_frame, mapping):
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if not run.text:
                continue
            for k, v in mapping.items():
                if k in run.text:
                    run.text = run.text.replace(k, v)


def replace_in_shape(shape, mapping):
    if shape.has_text_frame:
        replace_in_text_frame_preserve_font(shape.text_frame, mapping)

    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                replace_in_text_frame_preserve_font(cell.text_frame, mapping)

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shape.shapes:
            replace_in_shape(s, mapping)


def replace_placeholders(prs, mapping):
    for slide in prs.slides:
        for shape in slide.shapes:
            replace_in_shape(shape, mapping)

        if slide.has_notes_slide:
            for shape in slide.notes_slide.shapes:
                replace_in_shape(shape, mapping)

    for layout in prs.slide_layouts:
        for shape in layout.shapes:
            replace_in_shape(shape, mapping)

    for master in prs.slide_masters:
        for shape in master.shapes:
            replace_in_shape(shape, mapping)


# -------------------------------------------------
# MERGE PPTX
# -------------------------------------------------

UNIVERSIDAD_2QRS_MODEL_KEY = "UNIVERSIDAD_2QRS"


def _replace_rids_in_element(el, rid_map: dict[str, str]):
    """
    Reemplaza ids de relación dentro de un fragmento XML copiado.
    Esto es necesario para que imágenes y recursos relacionados apunten al slide destino.
    """
    for e in el.iter():
        if not hasattr(e, "attrib"):
            continue
        for attr_key, attr_val in list(e.attrib.items()):
            if attr_val in rid_map:
                e.attrib[attr_key] = rid_map[attr_val]


def _element_has_placeholder(el) -> bool:
    """
    Detecta si un elemento XML es un placeholder.
    Los placeholders heredados de layout/master no deben materializarse como shapes reales,
    porque pueden duplicar cajas de texto o alterar el diseño final.
    """
    for e in el.iter():
        if str(e.tag).lower().endswith("}ph"):
            return True
    return False


def _iter_copyable_sp_tree_children(sp_tree, *, skip_placeholders: bool):
    """
    Itera los hijos realmente copiables de un spTree.
    No copia propiedades internas del grupo raíz del slide/layout/master.
    """
    for child in list(sp_tree):
        tag = child.tag.lower()

        # Estos nodos son propiedades internas del grupo raíz.
        # Copiarlos genera XML inválido o comportamientos visuales extraños.
        if tag.endswith("nvgrpsppr") or tag.endswith("grpsppr"):
            continue

        if skip_placeholders and _element_has_placeholder(child):
            continue

        yield child


def _clone_child_with_relationships(dest_slide, src_part, child):
    """
    Clona un shape XML y recrea sus relaciones en el slide destino.
    Aplica especialmente a imágenes, fondos, logotipos y otros recursos embebidos.
    """
    cloned_child = deepcopy(child)

    rid_map = {}
    for rId, rel in src_part.rels.items():
        try:
            new_rId = dest_slide.part.relate_to(
                rel._target,
                rel.reltype,
                is_external=rel.is_external,
            )
            rid_map[rId] = new_rId
        except Exception:
            # Si una relación no aplica al shape clonado, se omite sin romper el PPTX.
            continue

    _replace_rids_in_element(cloned_child, rid_map)
    return cloned_child


def _copy_sp_tree_children_to_slide(dest_slide, src_part, src_sp_tree, *, skip_placeholders: bool):
    """
    Copia shapes desde un slide, layout o master hacia un slide destino.
    Se usa en UNIVERSIDAD_2QRS para materializar el diseño visual heredado
    antes de ocultar el master/layout del destino.
    """
    dest_sp_tree = dest_slide.shapes._spTree

    for child in _iter_copyable_sp_tree_children(
        src_sp_tree,
        skip_placeholders=skip_placeholders,
    ):
        cloned_child = _clone_child_with_relationships(dest_slide, src_part, child)
        dest_sp_tree.insert_element_before(cloned_child, "p:extLst")


def _remove_slide_level_shapes(slide):
    """
    Elimina shapes creados automáticamente al agregar una diapositiva nueva.
    No elimina elementos heredados del master, por eso en 2QRS también se usa
    showMasterSp=0 y una materialización controlada del diseño fuente.
    """
    for shape in list(slide.shapes):
        element = shape._element
        element.getparent().remove(element)


def _hide_inherited_master_shapes(slide):
    """
    Evita que la diapositiva clonada herede shapes del master del PPTX destino.

    Esto es clave para UNIVERSIDAD_2QRS: cuando se mezclan DIPLOMADO, CURSO,
    CERTIFICADO DE ESTUDIOS, etc. dentro del mismo modelo, usar el master/layout
    del primer template puede superponer o agrandar el fondo del jaguar/escudo.
    """
    slide._element.set("showMasterSp", "0")


def _count_non_placeholder_shapes_in_layout(layout) -> int:
    """
    Cuenta shapes visuales no-placeholder en un layout.
    Sirve para elegir el layout más limpio del destino en UNIVERSIDAD_2QRS.
    """
    count = 0
    for _ in _iter_copyable_sp_tree_children(
        layout.shapes._spTree,
        skip_placeholders=True,
    ):
        count += 1
    return count


def _get_cleanest_layout(prs: Presentation):
    """
    Devuelve el layout con menos shapes visuales propios.
    Es más seguro que asumir que slide_layouts[6] siempre está realmente vacío.
    """
    layouts = list(prs.slide_layouts)
    if not layouts:
        raise ValueError("La presentación no tiene layouts disponibles")

    return min(layouts, key=_count_non_placeholder_shapes_in_layout)


def clone_slide_into(dest_prs: Presentation, src_slide):
    """
    Clonación original para los modelos que ya funcionan correctamente.
    No se cambia su comportamiento para evitar afectar INSTITUTO, UNIVERSIDAD_AZUL,
    COLEGIO_ABOGADOS_CALLAO o COLEGIO_DE_PROFESORES_DEL_PERU.
    """
    blank_layout = dest_prs.slide_layouts[6]
    new_slide = dest_prs.slides.add_slide(blank_layout)

    spTree = new_slide.shapes._spTree
    src_spTree = src_slide.shapes._spTree

    for child in list(src_spTree):
        tag = child.tag.lower()

        if tag.endswith("nvgrpsppr") or tag.endswith("grpsppr"):
            continue

        spTree.insert_element_before(deepcopy(child), "p:extLst")

    rid_map = {}
    for rId, rel in src_slide.part.rels.items():
        try:
            new_rId = new_slide.part.relate_to(
                rel._target,
                rel.reltype,
                is_external=rel.is_external,
            )
            rid_map[rId] = new_rId
        except Exception:
            continue

    _replace_rids_in_element(new_slide._element, rid_map)


def clone_slide_into_universidad_2qrs(dest_prs: Presentation, src_slide):
    """
    Clonación especial SOLO para UNIVERSIDAD_2QRS.

    Motivo:
    - En 2QRS, algunas piezas visuales del diseño pueden vivir en master/layout.
    - Al mezclar distintos tipoModelo en el mismo modeloCertificado, el slide clonado
      puede heredar el master/layout del primer PPTX del lote.
    - Eso produce fondos gigantes, marcas de agua duplicadas o elementos superpuestos.

    Solución:
    1. Crear un slide con el layout más limpio del destino.
    2. Ocultar shapes heredados del master del destino.
    3. Materializar primero el master y layout del slide fuente.
    4. Copiar encima los shapes reales del slide fuente.
    """
    clean_layout = _get_cleanest_layout(dest_prs)
    new_slide = dest_prs.slides.add_slide(clean_layout)

    _hide_inherited_master_shapes(new_slide)
    _remove_slide_level_shapes(new_slide)

    # 1) Materializar shapes visuales del master fuente, sin placeholders.
    try:
        src_master = src_slide.slide_layout.slide_master
        _copy_sp_tree_children_to_slide(
            new_slide,
            src_master.part,
            src_master.shapes._spTree,
            skip_placeholders=True,
        )
    except Exception:
        pass

    # 2) Materializar shapes visuales del layout fuente, sin placeholders.
    try:
        src_layout = src_slide.slide_layout
        _copy_sp_tree_children_to_slide(
            new_slide,
            src_layout.part,
            src_layout.shapes._spTree,
            skip_placeholders=True,
        )
    except Exception:
        pass

    # 3) Copiar shapes propios del slide fuente, incluyendo placeholders ya rellenados.
    _copy_sp_tree_children_to_slide(
        new_slide,
        src_slide.part,
        src_slide.shapes._spTree,
        skip_placeholders=False,
    )

    return new_slide


def merge_presentations(
    presentations: List[Presentation],
    modelo_certificado: str = "",
) -> Presentation:
    if not presentations:
        raise ValueError("No hay presentaciones para unir")

    modelo_key = (modelo_certificado or "").upper().strip()

    # Usar la primera presentación como base preserva su tamaño de página,
    # tema y configuración general del archivo.
    dest = presentations[0]

    for prs in presentations[1:]:
        for slide in prs.slides:
            if modelo_key == UNIVERSIDAD_2QRS_MODEL_KEY:
                clone_slide_into_universidad_2qrs(dest, slide)
            else:
                clone_slide_into(dest, slide)

    return dest


# -------------------------------------------------
# GENERACIÓN DE PPTX POR ITEM  ✅ RESTAURADA
# -------------------------------------------------

def generar_presentacion_por_item(item: DiplomaRequest) -> Presentation:
    tipo = item.tipoModelo.upper().strip()
    modelo_cert = item.modeloCertificado.upper().strip()

    template_path = resolve_template_path(modelo_cert, tipo)

    if not os.path.exists(template_path):
        raise HTTPException(
            status_code=500,
            detail=f"No existe la plantilla: {template_path}"
        )

    prs = Presentation(template_path)
    
    # Determinar formato de fecha según modelo de certificado
    usar_fecha_larga = modelo_cert in MODELOS_FECHA_LARGA

    if usar_fecha_larga:
        fecha_inicio, fecha_fin = format_date_range_long_es(
            item.fechaInicio,
            item.fechaFin
        )
        fecha_emision_larga = format_date_long_es(item.fechaEmision)
        fecha_emision_corta = format_date_ddmmyyyy(item.fechaEmision)
        # fecha_emision = format_date_long_es(item.fechaEmision)
    else:
        fecha_inicio = format_date_ddmmyyyy(item.fechaInicio)
        fecha_fin = format_date_ddmmyyyy(item.fechaFin)
        fecha_emision_larga = format_date_long_es(item.fechaEmision)
        fecha_emision_corta = format_date_ddmmyyyy(item.fechaEmision)
        # fecha_emision = format_date_ddmmyyyy(item.fechaEmision)
    
    nombre_posterior = nombre_completo_capitalizado(item.nombres, item.apellidos)
    modulos = obtener_modulos_por_tema(tipo, item.temaDiplomado)
    
    ##horas de módulos correctamente distribuidas
    cantidad_modulos = MODULOS_COUNT[tipo]

    horas_por_modulo_global = calcular_horas_por_modulo(
        item.horasAcademicas,
        cantidad_modulos
    )
    
    horas_por_modulo = distribuir_horas_por_modulo(
        item.horasAcademicas,
        cantidad_modulos
    )

    mapping = {
        "{{MODELO}}": tipo,
        "{{MODELO_MINUSCULA}}": modelo_con_mayuscula_inicial(tipo),

        "{{NOMBRES}}": item.nombres.upper(),
        "{{APELLIDOS}}": item.apellidos.upper(),
        "{{TEMA_DIPLOMADO}}": item.temaDiplomado.upper(),

        "{{NOMBRE_COMPLETO_MINUSCULA}}": nombre_posterior,

        "{{FECHA_INICIO}}": fecha_inicio,
        "{{FECHA_FIN}}": fecha_fin,

        "{{HORAS_ACADEMICAS}}": format_two_digits_number(item.horasAcademicas),
        "{{CREDITOS_ACADEMICOS}}": format_two_digits_float(item.creditosAcademicos),
        
        "{{HORAS_MODULO}}": horas_por_modulo_global,

        "{{FOLIO_NUMERO}}": item.folioNumero,

        "{{FECHA_EMISION_LARGA}}": fecha_emision_larga,
        "{{FECHA_EMISION_CORTA}}": fecha_emision_corta,

        "{{CODE_STUDENT}}": (item.codigoEstudiante or "").strip().lower(),
        "{{INICIALES}}": obtener_iniciales_apellidos(item.apellidos),
        "{{CIUDAD}}": (item.ciudad or "").strip(),
    }

    for i, m in enumerate(modulos, start=1):
        mapping[f"{{{{MODULO_{i}}}}}"] = m
    
    ##Agregar horas por módulo al mapping       
    for i, h in enumerate(horas_por_modulo, start=1):
        mapping[f"{{{{HORAS_MODULO_{i}}}}}"] = str(h)


    replace_placeholders(prs, mapping)

    label_tabla = LABEL_TIPO_MAP.get(tipo)
    if label_tabla:
        ajustar_tabla_certificado_estudios_generico(
            prs,
            label=label_tabla,
            tema=item.temaDiplomado.upper(),
        )

    return prs




# -------------------------------------------------
# MERGE PPTX POR PAQUETE OPC - SOLO UNIVERSIDAD_2QRS
# -------------------------------------------------
#
# Motivo:
# El modelo UNIVERSIDAD_2QRS tiene elementos visuales sensibles en master/layout
# y/o relaciones internas de imágenes. Cuando se mezclan distintos tipoModelo
# dentro del mismo modeloCertificado, clonar slides con python-pptx puede hacer
# que las diapositivas 2, 3, 4... hereden el master/layout equivocado del primer
# PPTX del lote. Por eso se deforman imágenes/fondos, aunque el texto salga bien.
#
# Esta solución NO toca los otros modelos. Para UNIVERSIDAD_2QRS se evita clonar
# con python-pptx y se unen los PPTX a nivel de paquete .pptx/.zip, copiando cada
# slide junto con sus layouts, masters, themes, imágenes y relaciones.

CT_NS = "http://schemas.openxmlformats.org/package/2006/content-types"
REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

SLIDE_REL_TYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide"
SLIDE_MASTER_REL_TYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster"

ET.register_namespace("", CT_NS)
ET.register_namespace("p", P_NS)
ET.register_namespace("r", R_NS)


def _xml_bytes(root) -> bytes:
    return ET.tostring(root, encoding="UTF-8", xml_declaration=True)


def _read_zip_to_dict(pptx_bytes: bytes) -> dict[str, bytes]:
    with zipfile.ZipFile(BytesIO(pptx_bytes), "r") as z:
        return {name: z.read(name) for name in z.namelist()}


def _write_dict_to_zip(files: dict[str, bytes]) -> bytes:
    output = BytesIO()
    with zipfile.ZipFile(output, "w", compression=zipfile.ZIP_DEFLATED) as z:
        for name, data in files.items():
            z.writestr(name, data)
    return output.getvalue()


def _rels_path_for_part(part_path: str) -> str:
    folder, filename = posixpath.split(part_path)
    return posixpath.join(folder, "_rels", f"{filename}.rels")


def _normalize_internal_target(base_part_path: str, target: str) -> str:
    if target.startswith("/"):
        return posixpath.normpath(target.lstrip("/"))

    base_dir = posixpath.dirname(base_part_path)
    return posixpath.normpath(posixpath.join(base_dir, target))


def _relative_target(from_part_path: str, to_part_path: str) -> str:
    from_dir = posixpath.dirname(from_part_path)
    return posixpath.relpath(to_part_path, start=from_dir)


def _max_number_for_pattern(names: set[str], pattern: str) -> int:
    regex = re.compile(pattern)
    max_value = 0

    for name in names:
        match = regex.match(name)
        if match:
            max_value = max(max_value, int(match.group(1)))

    return max_value


def _init_part_counters(existing_names: set[str]) -> dict[str, int]:
    return {
        "slide": _max_number_for_pattern(existing_names, r"^ppt/slides/slide(\d+)\.xml$"),
        "layout": _max_number_for_pattern(existing_names, r"^ppt/slideLayouts/slideLayout(\d+)\.xml$"),
        "master": _max_number_for_pattern(existing_names, r"^ppt/slideMasters/slideMaster(\d+)\.xml$"),
        "theme": _max_number_for_pattern(existing_names, r"^ppt/theme/theme(\d+)\.xml$"),
        "chart": _max_number_for_pattern(existing_names, r"^ppt/charts/chart(\d+)\.xml$"),
        "media": _max_number_for_pattern(existing_names, r"^ppt/media/[A-Za-z_]*(\d+)\.[A-Za-z0-9]+$"),
        "embedding": _max_number_for_pattern(existing_names, r"^ppt/embeddings/embedding(\d+)\.[A-Za-z0-9]+$"),
        "ole": _max_number_for_pattern(existing_names, r"^ppt/embeddings/oleObject(\d+)\.bin$"),
        "diagram": _max_number_for_pattern(existing_names, r"^ppt/diagrams/[A-Za-z]+(\d+)\.xml$"),
    }


def _next_named_part(folder: str, prefix: str, ext: str, counter_key: str, existing_names: set[str], counters: dict[str, int]) -> str:
    while True:
        counters[counter_key] += 1
        candidate = f"{folder}/{prefix}{counters[counter_key]}{ext}"
        if candidate not in existing_names:
            existing_names.add(candidate)
            return candidate


def _allocate_new_part_name(src_part_path: str, existing_names: set[str], counters: dict[str, int]) -> str:
    folder, filename = posixpath.split(src_part_path)
    base, ext = posixpath.splitext(filename)

    if folder == "ppt/slides":
        return _next_named_part("ppt/slides", "slide", ".xml", "slide", existing_names, counters)

    if folder == "ppt/slideLayouts":
        return _next_named_part("ppt/slideLayouts", "slideLayout", ".xml", "layout", existing_names, counters)

    if folder == "ppt/slideMasters":
        return _next_named_part("ppt/slideMasters", "slideMaster", ".xml", "master", existing_names, counters)

    if folder == "ppt/theme":
        return _next_named_part("ppt/theme", "theme", ".xml", "theme", existing_names, counters)

    if folder == "ppt/charts":
        return _next_named_part("ppt/charts", "chart", ".xml", "chart", existing_names, counters)

    if folder == "ppt/media":
        # Conserva la extensión real de la imagen: .png, .jpg, .jpeg, .emf, etc.
        return _next_named_part("ppt/media", "image", ext, "media", existing_names, counters)

    if folder == "ppt/embeddings":
        if ext.lower() == ".bin":
            return _next_named_part("ppt/embeddings", "oleObject", ".bin", "ole", existing_names, counters)
        return _next_named_part("ppt/embeddings", "embedding", ext, "embedding", existing_names, counters)

    # Fallback seguro para partes menos comunes.
    if src_part_path not in existing_names:
        existing_names.add(src_part_path)
        return src_part_path

    i = 1
    while True:
        candidate = posixpath.join(folder, f"{base}_copy{i}{ext}")
        if candidate not in existing_names:
            existing_names.add(candidate)
            return candidate
        i += 1


def _content_type_maps(content_types_xml: bytes):
    root = ET.fromstring(content_types_xml)

    overrides: dict[str, str] = {}
    defaults: dict[str, str] = {}

    for child in root:
        tag = child.tag
        if tag.endswith("Override"):
            part_name = child.attrib.get("PartName", "")
            content_type = child.attrib.get("ContentType", "")
            if part_name and content_type:
                overrides[part_name.lstrip("/")] = content_type

        elif tag.endswith("Default"):
            extension = child.attrib.get("Extension", "")
            content_type = child.attrib.get("ContentType", "")
            if extension and content_type:
                defaults[extension.lower()] = content_type

    return root, overrides, defaults


def _ensure_content_type(
    dest_ct_root,
    dest_overrides: dict[str, str],
    dest_defaults: dict[str, str],
    src_overrides: dict[str, str],
    src_defaults: dict[str, str],
    src_part_path: str,
    new_part_path: str,
):
    src_content_type = src_overrides.get(src_part_path)

    if src_content_type:
        if new_part_path not in dest_overrides:
            ET.SubElement(
                dest_ct_root,
                f"{{{CT_NS}}}Override",
                {
                    "PartName": f"/{new_part_path}",
                    "ContentType": src_content_type,
                },
            )
            dest_overrides[new_part_path] = src_content_type
        return

    ext = posixpath.splitext(src_part_path)[1].lstrip(".").lower()
    if not ext:
        return

    src_default_type = src_defaults.get(ext)
    if src_default_type and ext not in dest_defaults:
        ET.SubElement(
            dest_ct_root,
            f"{{{CT_NS}}}Default",
            {
                "Extension": ext,
                "ContentType": src_default_type,
            },
        )
        dest_defaults[ext] = src_default_type


def _next_relationship_id(rels_root) -> str:
    max_id = 0

    for rel in rels_root:
        rid = rel.attrib.get("Id", "")
        match = re.match(r"rId(\d+)$", rid)
        if match:
            max_id = max(max_id, int(match.group(1)))

    return f"rId{max_id + 1}"


def _next_slide_id(presentation_root) -> str:
    max_id = 255

    for el in presentation_root.iter():
        if el.tag.endswith("sldId"):
            try:
                max_id = max(max_id, int(el.attrib.get("id", "0")))
            except ValueError:
                pass

    return str(max_id + 1)


def _next_master_id(presentation_root) -> str:
    max_id = 2147483647

    for el in presentation_root.iter():
        if el.tag.endswith("sldMasterId"):
            try:
                max_id = max(max_id, int(el.attrib.get("id", "0")))
            except ValueError:
                pass

    return str(max_id + 1)


def _find_or_create_sld_id_list(presentation_root):
    sld_id_list = presentation_root.find(f"{{{P_NS}}}sldIdLst")
    if sld_id_list is not None:
        return sld_id_list

    sld_id_list = ET.Element(f"{{{P_NS}}}sldIdLst")
    presentation_root.append(sld_id_list)
    return sld_id_list


def _find_or_create_master_id_list(presentation_root):
    master_id_list = presentation_root.find(f"{{{P_NS}}}sldMasterIdLst")
    if master_id_list is not None:
        return master_id_list

    master_id_list = ET.Element(f"{{{P_NS}}}sldMasterIdLst")

    # Ubicarlo antes de sldIdLst si existe, para mantener el orden normal del XML.
    children = list(presentation_root)
    sld_id_list = presentation_root.find(f"{{{P_NS}}}sldIdLst")

    if sld_id_list is not None:
        index = children.index(sld_id_list)
        presentation_root.insert(index, master_id_list)
    else:
        presentation_root.insert(0, master_id_list)

    return master_id_list


def _get_source_slide_parts(src_files: dict[str, bytes]) -> list[str]:
    presentation_xml = ET.fromstring(src_files["ppt/presentation.xml"])
    presentation_rels = ET.fromstring(src_files["ppt/_rels/presentation.xml.rels"])

    rel_by_id = {
        rel.attrib.get("Id"): rel
        for rel in presentation_rels
    }

    slide_parts: list[str] = []
    sld_id_list = presentation_xml.find(f"{{{P_NS}}}sldIdLst")

    if sld_id_list is None:
        return slide_parts

    for sld_id in sld_id_list.findall(f"{{{P_NS}}}sldId"):
        rid = sld_id.attrib.get(f"{{{R_NS}}}id")
        rel = rel_by_id.get(rid)

        if rel is None:
            continue

        if rel.attrib.get("Type") != SLIDE_REL_TYPE:
            continue

        target = rel.attrib.get("Target", "")
        slide_part = _normalize_internal_target("ppt/presentation.xml", target)

        if slide_part in src_files:
            slide_parts.append(slide_part)

    return slide_parts


def merge_pptx_packages_preserve_layouts(pptx_streams: list[bytes]) -> bytes:
    """
    Une PPTX preservando slides, layouts, masters, themes, imágenes y relaciones.

    Se usa SOLO para UNIVERSIDAD_2QRS porque allí el problema está en imágenes/fondos
    heredados al mezclar distintos tipoModelo. Este método evita esa herencia errónea
    porque copia las partes OPC reales de cada PPTX en lugar de reconstruir slides con
    el master/layout del primer archivo.
    """
    if not pptx_streams:
        raise ValueError("No hay PPTX para unir")

    # Primer PPTX como base: se mantiene intacto. Por eso tu primer certificado
    # siempre se veía bien; ahora los siguientes también conservarán sus partes.
    dest_files = _read_zip_to_dict(pptx_streams[0])
    existing_names = set(dest_files.keys())
    counters = _init_part_counters(existing_names)

    dest_ct_root, dest_overrides, dest_defaults = _content_type_maps(
        dest_files["[Content_Types].xml"]
    )

    dest_presentation_root = ET.fromstring(dest_files["ppt/presentation.xml"])
    dest_presentation_rels_root = ET.fromstring(dest_files["ppt/_rels/presentation.xml.rels"])

    dest_sld_id_list = _find_or_create_sld_id_list(dest_presentation_root)
    dest_master_id_list = _find_or_create_master_id_list(dest_presentation_root)

    added_master_parts: set[str] = set()

    def add_master_to_presentation(master_part_path: str):
        if master_part_path in added_master_parts:
            return

        new_rid = _next_relationship_id(dest_presentation_rels_root)
        target = _relative_target("ppt/presentation.xml", master_part_path)

        ET.SubElement(
            dest_presentation_rels_root,
            f"{{{REL_NS}}}Relationship",
            {
                "Id": new_rid,
                "Type": SLIDE_MASTER_REL_TYPE,
                "Target": target,
            },
        )

        ET.SubElement(
            dest_master_id_list,
            f"{{{P_NS}}}sldMasterId",
            {
                "id": _next_master_id(dest_presentation_root),
                f"{{{R_NS}}}id": new_rid,
            },
        )

        added_master_parts.add(master_part_path)

    for pptx_bytes in pptx_streams[1:]:
        src_files = _read_zip_to_dict(pptx_bytes)
        src_ct_root, src_overrides, src_defaults = _content_type_maps(
            src_files["[Content_Types].xml"]
        )

        part_map: dict[str, str] = {}
        new_master_parts_from_this_source: set[str] = set()

        def copy_part_recursive(src_part_path: str) -> str:
            if src_part_path in part_map:
                return part_map[src_part_path]

            if src_part_path not in src_files:
                return src_part_path

            new_part_path = _allocate_new_part_name(
                src_part_path,
                existing_names,
                counters,
            )
            part_map[src_part_path] = new_part_path

            dest_files[new_part_path] = src_files[src_part_path]

            _ensure_content_type(
                dest_ct_root,
                dest_overrides,
                dest_defaults,
                src_overrides,
                src_defaults,
                src_part_path,
                new_part_path,
            )

            if new_part_path.startswith("ppt/slideMasters/"):
                new_master_parts_from_this_source.add(new_part_path)

            src_rels_path = _rels_path_for_part(src_part_path)
            if src_rels_path in src_files:
                src_rels_root = ET.fromstring(src_files[src_rels_path])

                for rel in src_rels_root:
                    target_mode = rel.attrib.get("TargetMode", "")
                    if target_mode.lower() == "external":
                        continue

                    target = rel.attrib.get("Target", "")
                    if not target:
                        continue

                    target_part_path = _normalize_internal_target(src_part_path, target)
                    if target_part_path not in src_files:
                        continue

                    new_target_part_path = copy_part_recursive(target_part_path)
                    rel.attrib["Target"] = _relative_target(new_part_path, new_target_part_path)

                new_rels_path = _rels_path_for_part(new_part_path)
                dest_files[new_rels_path] = _xml_bytes(src_rels_root)

            return new_part_path

        source_slide_parts = _get_source_slide_parts(src_files)

        for source_slide_part in source_slide_parts:
            new_slide_part = copy_part_recursive(source_slide_part)

            new_rid = _next_relationship_id(dest_presentation_rels_root)
            slide_target = _relative_target("ppt/presentation.xml", new_slide_part)

            ET.SubElement(
                dest_presentation_rels_root,
                f"{{{REL_NS}}}Relationship",
                {
                    "Id": new_rid,
                    "Type": SLIDE_REL_TYPE,
                    "Target": slide_target,
                },
            )

            ET.SubElement(
                dest_sld_id_list,
                f"{{{P_NS}}}sldId",
                {
                    "id": _next_slide_id(dest_presentation_root),
                    f"{{{R_NS}}}id": new_rid,
                },
            )

        for master_part in sorted(new_master_parts_from_this_source):
            add_master_to_presentation(master_part)

    dest_files["ppt/presentation.xml"] = _xml_bytes(dest_presentation_root)
    dest_files["ppt/_rels/presentation.xml.rels"] = _xml_bytes(dest_presentation_rels_root)
    dest_files["[Content_Types].xml"] = _xml_bytes(dest_ct_root)

    return _write_dict_to_zip(dest_files)


def insertar_qr_en_presentacion_item(prs: Presentation, item: DiplomaRequest) -> None:
    """
    Inserta QR antes de guardar cada PPTX individual.
    Esto es especialmente importante para UNIVERSIDAD_2QRS, porque luego se hará
    merge por paquete y cada PPTX debe quedar autocontenido.
    """
    qr_url = build_qr_url(
        item.nombres,
        item.apellidos,
        item.temaDiplomado,
    )
    qr_image = generate_qr_image(qr_url)

    qr_size_cm = 2.0 if item.modeloCertificado.upper().strip() == "UNIVERSIDAD_AZUL" else 2.78

    for slide in prs.slides:
        insert_qr_at_placeholder_in_slide(
            slide,
            qr_image,
            qr_size_cm=qr_size_cm,
        )


def presentation_to_pptx_bytes(prs: Presentation) -> bytes:
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue()


def generar_pptx_universidad_2qrs_preservando_diseno(items: List[DiplomaRequest]) -> bytes:
    """
    Genera el PPTX final de UNIVERSIDAD_2QRS evitando el merge visual de python-pptx.
    """
    pptx_streams: list[bytes] = []

    for item in items:
        prs = generar_presentacion_por_item(item)
        insertar_qr_en_presentacion_item(prs, item)
        pptx_streams.append(presentation_to_pptx_bytes(prs))

    return merge_pptx_packages_preserve_layouts(pptx_streams)


# -------------------------------------------------
# ENDPOINTS
# -------------------------------------------------

def obtener_modelo_unico_del_lote(items: List[DiplomaRequest]) -> str:
    """
    Valida que el lote pertenezca a un solo modeloCertificado.

    Sí se permite mezclar distintos tipoModelo, por ejemplo:
    - UNIVERSIDAD_2QRS + DIPLOMADO
    - UNIVERSIDAD_2QRS + CURSO
    - UNIVERSIDAD_2QRS + CURSO_DE_CAPACITACION

    Lo que no debe ocurrir es mezclar UNIVERSIDAD_2QRS con UNIVERSIDAD_AZUL,
    INSTITUTO u otros modelos dentro del mismo PPTX.
    """
    modelo_base = items[0].modeloCertificado.upper().strip()

    for index, item in enumerate(items, start=1):
        modelo_actual = item.modeloCertificado.upper().strip()
        if modelo_actual != modelo_base:
            raise HTTPException(
                status_code=400,
                detail=(
                    "No mezcles distintos modeloCertificado en un mismo lote. "
                    f"El primer item usa {modelo_base}, pero el item {index} usa {modelo_actual}. "
                    "Puedes mezclar tipoModelo, pero no modeloCertificado."
                ),
            )

    return modelo_base


@app.get("/health")
def health():
    return {"status": "ok"}


@app.post("/api/diplomas")
def generate_pptx_batch(
    payload: BatchRequest,
):
    if not payload.items:
        raise HTTPException(status_code=400, detail="items no puede estar vacío")

    modelo_lote = obtener_modelo_unico_del_lote(payload.items)

    # CASO ESPECIAL: UNIVERSIDAD_2QRS
    # Este modelo se une a nivel de paquete .pptx/.zip para conservar exactamente
    # las imágenes, fondos, layouts y masters de cada tipoModelo.
    #
    # Esto corrige el problema donde el primer certificado salía bien, pero desde
    # el segundo se mezclaban/deformaban imágenes del certificado anterior.
    if modelo_lote == UNIVERSIDAD_2QRS_MODEL_KEY:
        merged_bytes = generar_pptx_universidad_2qrs_preservando_diseno(payload.items)
        output = BytesIO(merged_bytes)
        output.seek(0)

        filename = f"CERTIFICADOS_{int(datetime.now().timestamp())}.pptx"

        return StreamingResponse(
            output,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f'attachment; filename="{filename}"'}
        )

    # RESTO DE MODELOS:
    # Se mantiene el flujo anterior porque, según tus pruebas, no presentan el error.
    presentations: List[Presentation] = []
    slide_counts: List[int] = []

    for item in payload.items:
        prs = generar_presentacion_por_item(item)
        presentations.append(prs)
        slide_counts.append(len(prs.slides))

    merged = merge_presentations(
        presentations,
        modelo_certificado=modelo_lote,
    )

    merged_slides = list(merged.slides)
    cursor = 0

    for item, slide_count in zip(payload.items, slide_counts):
        qr_url = build_qr_url(
            item.nombres,
            item.apellidos,
            item.temaDiplomado,
        )
        qr_image = generate_qr_image(qr_url)

        qr_size_cm = 2.0 if item.modeloCertificado.upper().strip() == "UNIVERSIDAD_AZUL" else 2.78

        slides_del_item = merged_slides[cursor: cursor + slide_count]
        for slide in slides_del_item:
            insert_qr_at_placeholder_in_slide(
                slide,
                qr_image,
                qr_size_cm=qr_size_cm,
            )

        cursor += slide_count

    output = BytesIO()
    merged.save(output)
    output.seek(0)

    filename = f"CERTIFICADOS_{int(datetime.now().timestamp())}.pptx"

    return StreamingResponse(
        output,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'}
    )

